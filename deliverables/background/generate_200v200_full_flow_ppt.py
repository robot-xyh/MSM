#!/usr/bin/env python3
"""Generate the 200v200 full-flow engineering presentation."""

from __future__ import annotations

from pathlib import Path
import tempfile

from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
BACKGROUND = ROOT / "deliverables" / "background"
FULL = BACKGROUND / "figures" / "200v200_full_stack_3d"
LEARNING = BACKGROUND / "figures" / "200v200_learning_gnn"
AIRSIM = (
    ROOT
    / "research_modules"
    / "airsim_runtime"
    / "outputs"
    / "d5_cv_5v5_multicamera_formal_20260716"
    / "d5_backend_comparison.png"
)
OUTPUT = BACKGROUND / "200V200_三维全流程程序综合方案_汇报版.pptx"


SW = Inches(13.333)
SH = Inches(7.5)

FONT = "Microsoft YaHei"
FONT_MONO = "Consolas"

INK = "20282D"
MUTED = "5E6A70"
LIGHT_MUTED = "8A9499"
PAPER = "F4F6F7"
WHITE = "FFFFFF"
LINE = "D6DCDE"
TEAL = "147D76"
TEAL_LIGHT = "E4F1EF"
ORANGE = "D47A28"
ORANGE_LIGHT = "FBF0E3"
RED = "B84643"
RED_LIGHT = "F7E7E6"
BLUE = "2F6E9E"
BLUE_LIGHT = "E6EEF5"
PURPLE = "79558E"
PURPLE_LIGHT = "EEE8F2"
GREEN = "3F7D4B"
GREEN_LIGHT = "E8F1E9"
DARK = "242A2D"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run_font(run, size: float, color: str, *, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.04,
    font: str = FONT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(
    slide,
    lines: list[tuple[str, bool, str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 17,
    bullet: bool = False,
    gap: float = 7,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, (text, bold, color) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = ("• " if bullet else "") + text
        paragraph.level = 0
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.08
        for run in paragraph.runs:
            set_run_font(run, size, color, bold=bold)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = LINE,
    width: float = 1.0,
    radius: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    return shape


def add_label(
    slide,
    text: str,
    x: float,
    y: float,
    *,
    fill: str = TEAL_LIGHT,
    color: str = TEAL,
    width: float | None = None,
):
    w = width if width is not None else max(0.8, min(2.4, 0.16 * len(text) + 0.32))
    add_rect(slide, x, y, w, 0.34, fill=fill, line=fill, radius=True)
    add_text(
        slide,
        text,
        x + 0.05,
        y + 0.02,
        w - 0.10,
        0.27,
        size=10.5,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def add_footer(slide, page: int, *, dark: bool = False, note: str = "") -> None:
    color = "CBD2D5" if dark else LIGHT_MUTED
    add_text(
        slide,
        note or "200对200三维科研仿真方案",
        0.42,
        7.12,
        8.8,
        0.22,
        size=8.5,
        color=color,
        margin=0,
    )
    add_text(
        slide,
        f"{page:02d}",
        12.25,
        7.08,
        0.55,
        0.25,
        size=9,
        color=color,
        bold=True,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def add_title(
    slide,
    title: str,
    page: int,
    *,
    section: str = "",
    subtitle: str = "",
    dark: bool = False,
):
    background = DARK if dark else PAPER
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(background)
    if section:
        add_text(
            slide,
            section,
            0.45,
            0.22,
            2.1,
            0.28,
            size=10.5,
            color=ORANGE if dark else TEAL,
            bold=True,
            margin=0,
        )
    add_text(
        slide,
        title,
        0.45,
        0.53,
        11.9,
        0.55,
        size=25,
        color=WHITE if dark else INK,
        bold=True,
        margin=0,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            0.47,
            1.09,
            11.7,
            0.34,
            size=11.5,
            color="C8D0D3" if dark else MUTED,
            margin=0,
        )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.45),
        Inches(1.42),
        Inches(1.12),
        Inches(0.055),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(ORANGE)
    line.line.fill.background()
    add_footer(slide, page, dark=dark)


def add_picture_contain(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
):
    with Image.open(path) as image:
        iw, ih = image.size
    ratio = min(w / iw, h / ih)
    pw = iw * ratio
    ph = ih * ratio
    return slide.shapes.add_picture(
        str(path),
        Inches(x + (w - pw) / 2),
        Inches(y + (h - ph) / 2),
        Inches(pw),
        Inches(ph),
    )


def add_picture_crop(
    slide,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
):
    with Image.open(path) as image:
        iw, ih = image.size
    target = w / h
    source = iw / ih
    picture = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if source > target:
        visible = target / source
        picture.crop_left = (1.0 - visible) / 2
        picture.crop_right = (1.0 - visible) / 2
    else:
        visible = source / target
        picture.crop_top = (1.0 - visible) / 2
        picture.crop_bottom = (1.0 - visible) / 2
    return picture


def add_metric(
    slide,
    value: str,
    label: str,
    x: float,
    y: float,
    w: float,
    *,
    color: str = TEAL,
    note: str = "",
):
    add_rect(slide, x, y, w, 1.12, fill=WHITE, line=LINE)
    add_text(
        slide,
        value,
        x + 0.12,
        y + 0.11,
        w - 0.24,
        0.46,
        size=23,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        label,
        x + 0.10,
        y + 0.62,
        w - 0.20,
        0.24,
        size=10.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    if note:
        add_text(
            slide,
            note,
            x + 0.10,
            y + 0.88,
            w - 0.20,
            0.18,
            size=7.5,
            color=MUTED,
            align=PP_ALIGN.CENTER,
            margin=0,
        )


def add_table(
    slide,
    rows: list[list[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    widths: list[float] | None = None,
    font_size: float = 11.5,
):
    table = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)
    ).table
    if widths:
        total = sum(widths)
        for column, width in zip(table.columns, widths, strict=True):
            column.width = Inches(w * width / total)
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(
                DARK if row_index == 0 else (WHITE if row_index % 2 else "EDF1F2")
            )
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if column_index else PP_ALIGN.LEFT
                paragraph.space_after = Pt(0)
                for run in paragraph.runs:
                    set_run_font(
                        run,
                        font_size if row_index else font_size - 0.5,
                        WHITE if row_index == 0 else INK,
                        bold=row_index == 0,
                    )
    return table


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color=TEAL):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def add_flow_box(
    slide,
    title: str,
    body: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = WHITE,
    line: str = TEAL,
):
    add_rect(slide, x, y, w, h, fill=fill, line=line, width=1.3)
    add_text(
        slide,
        title,
        x + 0.12,
        y + 0.11,
        w - 0.24,
        0.29,
        size=13,
        color=line,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        body,
        x + 0.14,
        y + 0.46,
        w - 0.28,
        h - 0.54,
        size=10.5,
        color=INK,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def add_full_figure_slide(
    prs: Presentation,
    path: Path,
    page: int,
    *,
    tag: str,
    note: str,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)
    add_picture_contain(slide, path, 0.18, 0.08, 12.97, 6.92)
    add_label(slide, tag, 0.34, 0.24, fill=DARK, color=WHITE, width=1.25)
    add_footer(slide, page, note=note)
    return slide


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    prs.core_properties.title = "200对200三维全流程程序综合方案"
    prs.core_properties.subject = "科研仿真体系、模块原理、当前证据与实施路径"
    prs.core_properties.author = "MSM项目组"
    prs.core_properties.keywords = "反无人机, 多无人机, 三维仿真, 多传感器融合"

    # 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    with tempfile.NamedTemporaryFile(suffix=".png") as temporary:
        with Image.open(FULL / "02_layered_3d_scene.png").convert("RGB") as image:
            image = ImageEnhance.Brightness(image).enhance(0.32)
            image = ImageEnhance.Contrast(image).enhance(1.08)
            image.save(temporary.name)
        add_picture_crop(slide, Path(temporary.name), 0, 0, 13.333, 7.5)
    add_rect(slide, 0, 0, 5.72, 7.5, fill=DARK, line=DARK)
    add_text(
        slide,
        "200对200",
        0.62,
        0.88,
        4.6,
        0.65,
        size=34,
        color=ORANGE,
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "三维全流程程序综合方案",
        0.62,
        1.58,
        4.72,
        1.28,
        size=28,
        color=WHITE,
        bold=True,
        margin=0,
    )
    add_text(
        slide,
        "多传感器融合 · 轨迹关联 · 区域调度 · 资源分配\n降级接管 · 视觉配准 · 比例导引 · 系统评估",
        0.65,
        3.05,
        4.55,
        1.12,
        size=15.5,
        color="D8DEE0",
        margin=0,
    )
    add_text(
        slide,
        "科研仿真方案  |  2026年7月",
        0.65,
        6.58,
        4.4,
        0.30,
        size=11,
        color="B6C0C4",
        margin=0,
    )
    add_label(slide, "汇报版", 0.65, 5.84, fill=TEAL, color=WHITE, width=1.1)
    add_footer(slide, 1, dark=True, note="基于现有软件与仿真证据编制")

    # 2 Context
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "任务背景与设计目标", 2, section="01 体系定位")
    add_metric(slide, "200", "来袭目标", 0.55, 1.78, 2.15, color=RED)
    add_metric(slide, "200", "拦截资源", 2.93, 1.78, 2.15, color=TEAL)
    add_metric(slide, "8", "高空侦察节点", 5.31, 1.78, 2.15, color=ORANGE)
    add_metric(slide, "8", "区域数量", 7.69, 1.78, 2.15, color=BLUE)
    add_metric(slide, "3D", "北-东-地空间", 10.07, 1.78, 2.15, color=PURPLE)
    add_rect(slide, 0.55, 3.25, 5.95, 2.73, fill=WHITE, line=LINE)
    add_text(slide, "工程目标", 0.78, 3.52, 2.0, 0.34, size=17, color=TEAL, bold=True)
    add_rich_text(
        slide,
        [
            ("在统一时间轴和三维坐标系中形成全局航迹", False, INK),
            ("在目标与资源非等量条件下持续滚动分配", False, INK),
            ("中心失效或计划失配时保持权威连续和任务连续", False, INK),
            ("将雷达中段交接到末端视觉导引并保留安全门", False, INK),
        ],
        0.80,
        4.02,
        5.38,
        1.72,
        size=14.5,
        bullet=True,
        gap=8,
    )
    add_rect(slide, 6.78, 3.25, 5.95, 2.73, fill=WHITE, line=LINE)
    add_text(slide, "阶段目标", 7.02, 3.52, 2.0, 0.34, size=17, color=ORANGE, bold=True)
    add_rich_text(
        slide,
        [
            ("先冻结可解释、可回退的确定性基线", False, INK),
            ("再逐项验证区域策略、代价残差、跨视角图网络和主动视觉", False, INK),
            ("学习模型只给建议、概率或有界修正", False, INK),
            ("计划、身份、降级权威和飞行控制仍由确定性合同约束", False, INK),
        ],
        7.04,
        4.02,
        5.34,
        1.74,
        size=14.5,
        bullet=True,
        gap=8,
    )

    # 3 Evidence boundary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "证据边界", 3, section="01 体系定位", subtitle="所有结果按证据来源分级陈述")
    columns = [
        (
            "已验证",
            "软件单元测试\n三维质点仿真\nAirSim计算机视觉模式\nSimpleFlight受控仿真",
            GREEN,
            GREEN_LIGHT,
        ),
        (
            "机制说明",
            "算法流程图\n状态机与合同\n协方差、门控和求解关系\n不表示任务指标已达到",
            BLUE,
            BLUE_LIGHT,
        ),
        (
            "待验证",
            "真实传感器精度\n实时时延与通信容量\n硬件在环和实飞\n物理拦截与学习收益",
            ORANGE,
            ORANGE_LIGHT,
        ),
    ]
    for index, (title, body, color, fill) in enumerate(columns):
        x = 0.62 + index * 4.16
        add_rect(slide, x, 1.82, 3.73, 3.42, fill=fill, line=color, width=1.5)
        add_text(
            slide,
            title,
            x + 0.22,
            2.08,
            3.29,
            0.42,
            size=20,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            body,
            x + 0.34,
            2.78,
            3.05,
            1.86,
            size=15,
            color=INK,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    add_rect(slide, 0.62, 5.62, 12.0, 0.82, fill=RED_LIGHT, line=RED)
    add_text(
        slide,
        "当前没有实飞和硬件在环证据；期望效果图不能解释为实测增益。",
        0.90,
        5.85,
        11.44,
        0.34,
        size=16,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )

    # 4 Overall judgement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "总体技术判断", 4, section="01 体系定位")
    add_rect(slide, 0.55, 1.73, 12.2, 1.08, fill=TEAL_LIGHT, line=TEAL)
    add_text(
        slide,
        "200阶线性指派本身可计算，主要瓶颈在候选生成、代价构造、对象物化和证据写出。",
        0.83,
        2.03,
        11.65,
        0.42,
        size=18,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_metric(slide, "4万项", "200×200完整矩阵", 0.72, 3.20, 2.58, color=BLUE)
    add_metric(slide, "6400边", "稀疏候选基准", 3.56, 3.20, 2.58, color=TEAL)
    add_metric(slide, "85.4 ms", "向量化独立基准", 6.40, 3.20, 2.58, color=GREEN)
    add_metric(slide, "1904.3 ms", "旧逐边Python路径", 9.24, 3.20, 2.58, color=RED)
    add_flow_box(
        slide,
        "确定性基线",
        "扩展卡尔曼滤波、稀疏门控、匈牙利、规则区域压力、版本化状态机、几何配准、比例导引",
        0.72,
        4.83,
        5.67,
        1.35,
        fill=WHITE,
        line=TEAL,
    )
    add_flow_box(
        slide,
        "学习增强",
        "区域图策略、D3有界代价残差、D5跨视角图网络、主动视觉策略；默认关闭并保留规则回退",
        6.93,
        4.83,
        5.67,
        1.35,
        fill=WHITE,
        line=PURPLE,
    )

    # 5 Architecture
    add_full_figure_slide(
        prs,
        FULL / "01_full_stack_architecture.png",
        5,
        tag="总体架构",
        note="算法机制图：确定性在线主线与默认关闭的学习增强层",
    )

    # 6 Scene
    add_full_figure_slide(
        prs,
        FULL / "02_layered_3d_scene.png",
        6,
        tag="质点仿真",
        note="实际离线回合轨迹：seed 1000；右侧分层职责为方案说明",
    )

    # 7 Parameters
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "名义三维场景", 7, section="02 场景与合同")
    rows = [
        ["参数", "当前仿真值", "作用"],
        ["世界半边长 / 高度", "6000 m / 30–1500 m", "目标与资源生成边界"],
        ["目标 / 拦截器速度", "3.5–4.7 / 14 m·s⁻¹", "名义运动与导航上限"],
        ["主相机", "1920×1080，90°", "拦截资源末端视觉"],
        ["侦察相机", "3840×2160，70°", "高空侦察和跨视角锚点"],
        ["雷达 / 视觉 / 声学周期", "0.2 / 0.1 / 0.5 s", "异步观测节拍"],
        ["D3 / D4周期", "1.0 / 5.0 s", "目标分配 / 区域建议"],
        ["物理步长 / 判定半径", "0.05 s / 5 m", "状态推进 / 离线结果"],
    ]
    add_table(slide, rows, 0.62, 1.72, 8.15, 4.85, widths=[2.5, 2.4, 3.1], font_size=11.2)
    add_rect(slide, 9.08, 1.72, 3.65, 4.85, fill=WHITE, line=LINE)
    add_text(slide, "空间组织", 9.36, 2.03, 3.08, 0.35, size=17, color=TEAL, bold=True)
    add_rich_text(
        slide,
        [
            ("北-东-地（NED）六维状态", False, INK),
            ("水平面固定划分8个扇区", False, INK),
            ("每区维护需求、威胁、资源、通信和权威状态", False, INK),
            ("真实设备指标需由硬件试验重新辨识", True, RED),
        ],
        9.34,
        2.60,
        3.05,
        2.50,
        size=14,
        bullet=True,
        gap=10,
    )
    add_label(slide, "仿真参数", 9.35, 5.64, fill=ORANGE_LIGHT, color=ORANGE, width=1.28)

    # 8 Contracts
    add_full_figure_slide(
        prs,
        FULL / "03_contract_chain.png",
        8,
        tag="统一合同",
        note="双时间戳、协方差、NED状态、身份权威、计划版本和离线真值隔离",
    )

    # 9 Episode
    add_full_figure_slide(
        prs,
        FULL / "11_episode_state_machine.png",
        9,
        tag="MAIN",
        note="主运行时统一场景、时钟、通信、模块调度、执行确认和离线评估入口",
    )

    # 10 D1 overview
    add_full_figure_slide(
        prs,
        FULL / "04_d1_sensor_fusion.png",
        10,
        tag="D1",
        note="算法机制图：异步雷达、声学和光电观测统一为带协方差的GlobalTrack",
    )

    # 11 D1 mechanics
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "D1异步融合与乱序量测", 11, section="03 感知与关联")
    stages = [
        ("观测检查", "双时间戳\n协方差\n坐标与来源"),
        ("量测时刻更新", "雷达球坐标\n声学方位\n光电像素"),
        ("固定滞后回放", "6秒历史窗\n0.1秒时间桶\n重复谱系去重"),
        ("发布时刻补偿", "状态前向传播\n协方差增长\n质量分级"),
    ]
    for index, (title, body) in enumerate(stages):
        x = 0.58 + index * 3.12
        add_flow_box(
            slide,
            title,
            body,
            x,
            2.03,
            2.62,
            1.72,
            fill=WHITE,
            line=[BLUE, TEAL, ORANGE, PURPLE][index],
        )
        if index < 3:
            add_arrow(slide, x + 2.64, 2.88, x + 3.02, 2.88)
    add_rect(slide, 0.58, 4.23, 12.15, 1.55, fill=WHITE, line=LINE)
    add_text(slide, "协方差处理", 0.85, 4.52, 1.75, 0.32, size=16, color=TEAL, bold=True)
    add_text(
        slide,
        "距离、信噪比、遮挡、检测框尺度、量测延迟和外参误差进入量测协方差。\n"
        "过程噪声吸收转弯和加速；归一化创新平方用于在线一致性检查。\n"
        "位置误差等真值指标仅由D6在回合结束后离线计算。",
        2.28,
        4.46,
        9.95,
        0.84,
        size=14.5,
        color=INK,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_label(slide, "默认：扩展卡尔曼滤波", 0.85, 5.28, fill=BLUE_LIGHT, color=BLUE, width=2.08)
    add_label(slide, "预留：交互多模型 / 无迹滤波", 3.15, 5.28, fill=PURPLE_LIGHT, color=PURPLE, width=2.52)

    # 12 D1 results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "D1三维融合结果", 12, section="03 感知与关联", subtitle="seed 1000，10秒质点仿真")
    add_picture_contain(slide, FULL / "12_d1_3d_fusion_experiment.png", 0.48, 1.55, 8.25, 4.95)
    add_metric(slide, "201", "末帧融合航迹", 8.95, 1.72, 1.72, color=TEAL)
    add_metric(slide, "200", "离线几何配对", 10.86, 1.72, 1.72, color=BLUE)
    add_metric(slide, "6.56 m", "配对距离中位数", 8.95, 3.13, 1.72, color=GREEN)
    add_metric(slide, "24.92 m", "配对距离P95", 10.86, 3.13, 1.72, color=ORANGE)
    add_metric(slide, "1.0479", "平均归一化创新平方", 8.95, 4.54, 1.72, color=PURPLE)
    add_metric(slide, "99.29%", "创新门内比例", 10.86, 4.54, 1.72, color=TEAL)
    add_text(
        slide,
        "11.36 m为末帧几何配对均方根，不是严格身份连续的正式位置均方根误差。",
        8.96,
        6.03,
        3.62,
        0.42,
        size=10.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 13 D2 overview
    add_full_figure_slide(
        prs,
        FULL / "05_d2_data_association.png",
        13,
        tag="D2",
        note="算法机制图：稀疏候选、马氏距离门、全局最近邻和稳定航迹编号",
    )

    # 14 D2 lifecycle
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "D2关联与生命周期", 14, section="03 感知与关联")
    life = [
        ("暂定", "连续命中与风险检查"),
        ("确认", "稳定状态与质量"),
        ("可执行", "满足下游使用条件"),
        ("丢失", "有界预测和等待恢复"),
        ("删除", "超时或质量不可恢复"),
    ]
    for index, (title, body) in enumerate(life):
        x = 0.55 + index * 2.48
        add_flow_box(
            slide,
            title,
            body,
            x,
            1.83,
            2.05,
            1.38,
            fill=[BLUE_LIGHT, TEAL_LIGHT, GREEN_LIGHT, ORANGE_LIGHT, RED_LIGHT][index],
            line=[BLUE, TEAL, GREEN, ORANGE, RED][index],
        )
        if index < len(life) - 1:
            add_arrow(slide, x + 2.06, 2.52, x + 2.38, 2.52, color=MUTED)
    add_rect(slide, 0.55, 3.68, 5.94, 2.16, fill=WHITE, line=LINE)
    add_text(slide, "默认关联", 0.82, 3.96, 1.7, 0.32, size=17, color=TEAL, bold=True)
    add_rich_text(
        slide,
        [
            ("空间索引与时间窗生成稀疏候选", False, INK),
            ("马氏距离删除不可行边", False, INK),
            ("运动、类别、来源和历史代价排序", False, INK),
            ("匈牙利算法完成一对一全局匹配", False, INK),
        ],
        0.82,
        4.43,
        5.30,
        1.12,
        size=13.5,
        bullet=True,
        gap=6,
    )
    add_rect(slide, 6.79, 3.68, 5.94, 2.16, fill=WHITE, line=LINE)
    add_text(slide, "升级条件", 7.06, 3.96, 1.7, 0.32, size=17, color=PURPLE, bold=True)
    add_rich_text(
        slide,
        [
            ("协方差椭球持续重叠", False, INK),
            ("单帧歧义无法在下一帧消除", False, INK),
            ("身份切换成为主要任务损失", False, INK),
            ("计算预算允许联合概率或多假设对照", False, INK),
        ],
        7.06,
        4.43,
        5.30,
        1.12,
        size=13.5,
        bullet=True,
        gap=6,
    )
    add_text(
        slide,
        "D2中的GNN指“全局最近邻”，与D5的图神经网络不是同一算法。",
        0.73,
        6.20,
        11.9,
        0.32,
        size=12,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 15 D2 result
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "D2三维关联结果", 15, section="03 感知与关联", subtitle="seed 1000，48个关联帧")
    add_picture_contain(slide, FULL / "13_d2_3d_association_experiment.png", 0.48, 1.55, 8.45, 4.95)
    add_metric(slide, "201", "末帧航迹", 9.14, 1.72, 1.58, color=TEAL)
    add_metric(slide, "200 / 1", "可用 / 歧义映射", 10.92, 1.72, 1.58, color=ORANGE)
    add_metric(slide, "99.81%", "可评分映射覆盖", 9.14, 3.15, 1.58, color=GREEN)
    add_metric(slide, "72.92%", "完整帧覆盖", 10.92, 3.15, 1.58, color=BLUE)
    add_metric(slide, "48", "身份切换下界", 9.14, 4.58, 1.58, color=RED)
    add_metric(slide, "18", "多真实目标候选映射", 10.92, 4.58, 1.58, color=PURPLE)
    add_text(
        slide,
        "严格身份切换、连续性和重复航迹指标仍不可用。\n"
        "当前优先减少融合谱系污染，并补齐全时序标签。",
        9.10,
        6.03,
        3.48,
        0.47,
        size=10.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 16 Regional deterministic
    add_full_figure_slide(
        prs,
        LEARNING / "08_deterministic_regional_flow.png",
        16,
        tag="区域层",
        note="确定性区域调度：压力、预测需求、资源守恒、备用和跨区容量",
    )

    # 17 Regional 3D
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "区域资源调度", 17, section="04 调度与分配")
    add_picture_contain(slide, LEARNING / "09_deterministic_regional_3d.png", 0.45, 1.55, 7.75, 5.10)
    add_rect(slide, 8.48, 1.70, 4.26, 4.78, fill=WHITE, line=LINE)
    add_text(slide, "两级确定性求解", 8.78, 1.98, 3.68, 0.36, size=18, color=TEAL, bold=True)
    add_rich_text(
        slide,
        [
            ("预测未来窗口内各区域目标需求", False, INK),
            ("形成配额、最低保障和备用量", False, INK),
            ("最小费用流决定区域间转移数量", False, INK),
            ("线性指派选择具体调动资源", False, INK),
            ("D3在目标需求槽上完成区域内分配", False, INK),
            ("计划迟滞、版本、租约和执行确认", False, INK),
        ],
        8.77,
        2.58,
        3.62,
        2.68,
        size=13.5,
        bullet=True,
        gap=8,
    )
    add_label(slide, "默认：规则贪心回退", 8.79, 5.54, fill=TEAL_LIGHT, color=TEAL, width=1.94)
    add_label(slide, "扩展：最小费用流", 10.87, 5.54, fill=BLUE_LIGHT, color=BLUE, width=1.62)

    # 18 D3 core
    add_full_figure_slide(
        prs,
        FULL / "06_d3_assignment_and_residual.png",
        18,
        tag="D3",
        note="确定性规则代价、需求槽、稀疏候选、匈牙利和可选有界残差",
    )

    # 19 D3 M-to-N
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "D3多资源任务与计划稳定", 19, section="04 调度与分配")
    add_flow_box(
        slide,
        "高威胁目标",
        "目标T001\n需求：2主用 + 1备用",
        0.62,
        1.86,
        2.30,
        1.42,
        fill=RED_LIGHT,
        line=RED,
    )
    slots = [
        ("主用槽1", "wave 0\n独立门控", TEAL),
        ("主用槽2", "wave 0\n独立门控", BLUE),
        ("备用槽", "wave 1\n待命", ORANGE),
    ]
    for index, (title, body, color) in enumerate(slots):
        x = 4.10 + index * 2.70
        add_flow_box(slide, title, body, x, 1.86, 2.20, 1.42, fill=WHITE, line=color)
        add_arrow(slide, 2.93, 2.57, x - 0.10, 2.57, color=color)
    add_rect(slide, 0.62, 3.75, 5.85, 2.18, fill=WHITE, line=LINE)
    add_text(slide, "全有或全无准入", 0.90, 4.05, 2.4, 0.32, size=17, color=TEAL, bold=True)
    add_text(
        slide,
        "必要槽只获得部分资源时，低优先级不完整联盟整体退出并重新求解。\n"
        "备用资源只有新版本显式激活后才能执行。\n"
        "避免“需要3架、实际1架”被记为满足。",
        0.90,
        4.50,
        5.18,
        0.98,
        size=14.2,
        color=INK,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_rect(slide, 6.79, 3.75, 5.94, 2.18, fill=WHITE, line=LINE)
    add_text(slide, "滚动规划与迟滞", 7.07, 4.05, 2.4, 0.32, size=17, color=ORANGE, bold=True)
    add_text(
        slide,
        "普通换配要求新代价至少改善20%，且驻留时间达到2秒。\n"
        "资源失效、需求变化、旧边不可行或高威胁未满足时可绕过普通迟滞。\n"
        "计划标识和版本必须严格连续。",
        7.07,
        4.50,
        5.20,
        0.98,
        size=14.2,
        color=INK,
        valign=MSO_ANCHOR.MIDDLE,
    )

    # 20 D3 result
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "D3 200对200分配结果", 20, section="04 调度与分配", subtitle="seed 1000，第8版计划")
    add_picture_contain(slide, FULL / "14_d3_3d_assignment_experiment.png", 0.45, 1.55, 8.45, 5.02)
    add_metric(slide, "200 / 200", "唯一资源 / 唯一航迹", 9.12, 1.70, 3.40, color=TEAL)
    add_metric(slide, "10", "计划确认", 9.12, 3.10, 1.58, color=BLUE)
    add_metric(slide, "1978", "绑定与控制采用", 10.94, 3.10, 1.58, color=GREEN)
    add_metric(slide, "221.7 ms", "D3集成平均耗时", 9.12, 4.52, 1.58, color=ORANGE)
    add_metric(slide, "284.3 ms", "D3集成P95", 10.94, 4.52, 1.58, color=RED)
    add_text(
        slide,
        "中位初始距离约3114米；10秒内无五米接近。\n"
        "该页证明分配与版本消费，不证明导引成功。",
        9.10,
        5.96,
        3.46,
        0.53,
        size=10.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 21 D3 residual
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "D3有界代价残差", 21, section="05 学习增强", subtitle="强化学习修正排序，确定性求解保持不变")
    add_picture_contain(slide, LEARNING / "04_d3_cost_residual_assignment.png", 0.55, 1.62, 8.05, 4.45)
    add_rect(slide, 8.88, 1.70, 3.85, 4.64, fill=WHITE, line=PURPLE)
    add_text(
        slide,
        "C_final = C_rule + 0.25·tanh(a)",
        9.15,
        2.03,
        3.30,
        0.46,
        size=18,
        color=PURPLE,
        bold=True,
        align=PP_ALIGN.CENTER,
        font=FONT_MONO,
    )
    add_rich_text(
        slide,
        [
            ("只对通过硬门的稀疏候选边修正", False, INK),
            ("不能生成不可达边或改写航迹身份", False, INK),
            ("同一匈牙利与需求槽求解器完成分配", False, INK),
            ("20/20改变代价，0/20改变正式绑定", True, RED),
            ("当前未证明任务收益", True, RED),
        ],
        9.18,
        2.82,
        3.22,
        2.55,
        size=13.5,
        bullet=True,
        gap=9,
    )
    add_label(slide, "影子模式", 9.18, 5.73, fill=PURPLE_LIGHT, color=PURPLE, width=1.24)
    add_label(slide, "默认关闭", 10.60, 5.73, fill=RED_LIGHT, color=RED, width=1.24)

    # 22 D4 figure
    add_full_figure_slide(
        prs,
        FULL / "07_d4_degradation_state_machine.png",
        22,
        tag="D4",
        note="主动风险仲裁、中心重规划、机动高空侦察二级节点和完全分布式保底",
    )

    # 23 D4 details
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "D4降级判定与二级节点", 23, section="06 降级与末端")
    blocks = [
        (
            "中心正常",
            "中心计划为当前版本\n心跳、时期和摘要可信",
            GREEN,
            GREEN_LIGHT,
        ),
        (
            "二级辅助",
            "补充图像和跨视角线索\n不转移计划owner",
            BLUE,
            BLUE_LIGHT,
        ),
        (
            "二级接管",
            "覆盖、注册、新鲜度、通信、租约和版本全部通过",
            PURPLE,
            PURPLE_LIGHT,
        ),
        (
            "完全分布式",
            "中心与二级均不可用\n联盟成员确认后原子提交",
            ORANGE,
            ORANGE_LIGHT,
        ),
    ]
    for index, (title, body, color, fill) in enumerate(blocks):
        x = 0.55 + index * 3.12
        add_flow_box(slide, title, body, x, 1.75, 2.64, 1.62, fill=fill, line=color)
        if index < 3:
            add_arrow(slide, x + 2.66, 2.56, x + 3.02, 2.56, color=MUTED)
    add_rect(slide, 0.55, 3.82, 5.90, 2.22, fill=WHITE, line=LINE)
    add_text(slide, "主动降级证据", 0.82, 4.10, 2.2, 0.32, size=17, color=TEAL, bold=True)
    add_rich_text(
        slide,
        [
            ("D1协方差与量测年龄", False, INK),
            ("D2歧义、身份切换与重复航迹", False, INK),
            ("D3计划年龄、可行性和代价裕度", False, INK),
            ("D5错配、友方冲突和末端适用性", False, INK),
        ],
        0.82,
        4.54,
        5.20,
        1.20,
        size=13.5,
        bullet=True,
        gap=6,
    )
    add_rect(slide, 6.78, 3.82, 5.95, 2.22, fill=WHITE, line=LINE)
    add_text(slide, "当前AirSim现象", 7.06, 4.10, 2.2, 0.32, size=17, color=ORANGE, bold=True)
    add_rich_text(
        slide,
        [
            ("正常中心跨视角关联：4", False, INK),
            ("二级接管与完全分布式：0", False, INK),
            ("网络联合覆盖均值：约0.65–0.72", False, INK),
            ("主要断点：并非所有目标同时可见", True, RED),
        ],
        7.06,
        4.54,
        5.20,
        1.20,
        size=13.5,
        bullet=True,
        gap=6,
    )

    # 24 D5 geometry
    add_full_figure_slide(
        prs,
        FULL / "08_d5_geometry_registration.png",
        24,
        tag="D5",
        note="算法机制图：中心三维航迹预测、相机投影、协方差门控和跨视角绑定",
    )

    # 25 D5 controlled test
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "D5多相机受控结果", 25, section="06 降级与末端", subtitle="5主相机 + 1侦察相机，5目标，12秒，seed 7")
    add_picture_contain(slide, AIRSIM, 0.52, 1.52, 7.10, 4.20)
    rows = [
        ["后端", "检测召回", "严格配准", "稳定配准", "身份切换"],
        ["AirSim detect", "1.000", "1.000", "0.975", "0"],
        ["YOLOv8 + ByteTrack", "0.622", "0.966", "0.955", "25"],
    ]
    add_table(slide, rows, 7.82, 1.73, 4.92, 1.78, widths=[2.3, 1.1, 1.1, 1.1, 1.0], font_size=9.8)
    add_rect(slide, 7.82, 3.79, 4.92, 1.76, fill=WHITE, line=LINE)
    add_text(slide, "闭合路径", 8.10, 4.06, 1.30, 0.30, size=16, color=TEAL, bold=True)
    add_text(
        slide,
        "局部子集视场 → 侦察锚点 → 中心投影\n"
        "→ 匈牙利匹配 → 稳定窗口",
        8.10,
        4.50,
        4.35,
        0.60,
        size=14,
        color=INK,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_rect(slide, 7.82, 5.78, 4.92, 0.70, fill=RED_LIGHT, line=RED)
    add_text(
        slide,
        "该专项由main合成中心航迹，不代表真实雷达到视觉链已闭合。",
        8.04,
        5.97,
        4.48,
        0.29,
        size=10.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )

    # 26 D5 GNN
    add_full_figure_slide(
        prs,
        LEARNING / "03_gnn_crossview_matching.png",
        26,
        tag="D5学习",
        note="图神经网络只输出匿名轨迹同目标概率；几何候选、聚类约束和中心绑定保持确定性",
    )

    # 27 Active vision
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "侦察主动视觉", 27, section="06 降级与末端")
    add_picture_contain(slide, LEARNING / "06_active_vision_reconnaissance.png", 0.50, 1.55, 7.35, 4.95)
    add_rect(slide, 8.12, 1.68, 4.60, 2.12, fill=WHITE, line=TEAL)
    add_text(slide, "确定性默认路径", 8.42, 1.95, 3.98, 0.34, size=18, color=TEAL, bold=True)
    add_rich_text(
        slide,
        [
            ("已分配目标优先观察", False, INK),
            ("短时丢失进入重捕获并切广角", False, INK),
            ("无目标时确定性扫描扇区", False, INK),
            ("云台限速、限位、版本和友方预留复核", False, INK),
        ],
        8.40,
        2.46,
        3.95,
        1.12,
        size=12.8,
        bullet=True,
        gap=5,
    )
    add_rect(slide, 8.12, 4.02, 4.60, 2.10, fill=WHITE, line=PURPLE)
    add_text(slide, "学习候选现状", 8.42, 4.29, 3.98, 0.34, size=18, color=PURPLE, bold=True)
    add_rich_text(
        slide,
        [
            ("115万样本，总体准确率约0.956", False, INK),
            ("重捕获占92.16%，类别明显失衡", False, INK),
            ("观察目标召回为0，保持动作无正样本", True, RED),
            ("辅助和相机执行权限保持关闭", True, RED),
        ],
        8.40,
        4.80,
        3.95,
        1.12,
        size=12.8,
        bullet=True,
        gap=5,
    )

    # 28 D7
    add_full_figure_slide(
        prs,
        FULL / "09_d7_guidance_handover.png",
        28,
        tag="D7",
        note="位置比例导引到视觉比例导引的合同化交接；距离只打开尝试窗口",
    )

    # 29 D6
    add_full_figure_slide(
        prs,
        FULL / "10_d6_evaluation_and_timing.png",
        29,
        tag="D6",
        note="离线真值联接、指标可用性、逐seed结果、置信区间和阶段耗时",
    )

    # 30 Runtime timing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "200规模运行性能", 30, section="07 评估与实施", subtitle="seed 1000单回合；不同模块调用频率不同")
    timing = [
        ("D1融合", 29.1, TEAL),
        ("D2关联", 70.0, BLUE),
        ("D3分配", 221.7, ORANGE),
        ("D4区域建议", 2.9, GREEN),
        ("D5终端配准", 14.6, PURPLE),
        ("D5主动视觉", 38.7, RED),
        ("D7导引", 30.7, TEAL),
    ]
    max_value = max(value for _, value, _ in timing)
    for index, (name, value, color) in enumerate(timing):
        y = 1.78 + index * 0.62
        add_text(slide, name, 0.70, y + 0.05, 1.38, 0.28, size=12.5, color=INK, bold=True, align=PP_ALIGN.RIGHT)
        add_rect(slide, 2.28, y, 7.00, 0.38, fill="E2E7E9", line="E2E7E9")
        width = 7.00 * value / max_value
        add_rect(slide, 2.28, y, width, 0.38, fill=color, line=color)
        add_text(slide, f"{value:.1f} ms", 9.52, y + 0.03, 1.05, 0.28, size=11.5, color=color, bold=True, margin=0)
    add_rect(slide, 10.77, 1.78, 1.88, 2.20, fill=RED_LIGHT, line=RED)
    add_text(slide, "0.1247", 10.96, 2.14, 1.48, 0.55, size=25, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "实时因子", 10.96, 2.78, 1.48, 0.31, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "80.198秒墙钟\n对应10秒仿真", 10.98, 3.18, 1.44, 0.54, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 10.77, 4.35, 1.88, 1.65, fill=TEAL_LIGHT, line=TEAL)
    add_text(slide, "85.4 ms", 10.95, 4.70, 1.50, 0.43, size=20, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "D3独立向量化基准", 10.95, 5.22, 1.50, 0.48, size=10.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "优先优化D1批量数值、D2候选生成、D3对象与证据路径、日志写盘；不能用删除审计字段换取表面实时。",
        0.74,
        6.34,
        11.86,
        0.37,
        size=12.2,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 31 Safety shell
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "学习增强的确定性安全外壳", 31, section="07 评估与实施")
    layers = [
        ("输入检查", "双时间戳、协方差、版本、模型摘要、真值隔离", BLUE, BLUE_LIGHT),
        ("候选检查", "不可达、友方冲突、非邻接、过期和几何不一致先删除", TEAL, TEAL_LIGHT),
        ("动作投影", "残差幅度、资源守恒、边容量、云台限位和节点互斥", ORANGE, ORANGE_LIGHT),
        ("权威复核", "全局航迹、计划版本、owner、代次、租约、联盟和控制合同", PURPLE, PURPLE_LIGHT),
        ("规则回退", "模型缺失、超时、低置信、分布外或非法输出立即回退", RED, RED_LIGHT),
    ]
    widths = [11.60, 10.65, 9.70, 8.75, 7.80]
    for index, ((title, body, color, fill), width) in enumerate(zip(layers, widths, strict=True)):
        x = (13.333 - width) / 2
        y = 1.68 + index * 0.91
        add_rect(slide, x, y, width, 0.70, fill=fill, line=color, width=1.2)
        add_text(slide, title, x + 0.18, y + 0.16, 1.65, 0.29, size=14, color=color, bold=True, margin=0)
        add_text(slide, body, x + 1.92, y + 0.14, width - 2.10, 0.34, size=12.5, color=INK, align=PP_ALIGN.CENTER, margin=0)
    add_text(
        slide,
        "学习模型不能取得身份、计划发布、降级所有权或飞行控制权限。",
        1.18,
        6.40,
        11.0,
        0.34,
        size=15,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 32 Evidence learning
    add_full_figure_slide(
        prs,
        LEARNING / "07_current_evidence_boundary.png",
        32,
        tag="学习证据",
        note="数据、模型和影子接口已形成；在线辅助、身份、计划和控制权限仍关闭",
    )

    # 33 Experiment design
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "多随机种子实验设计", 33, section="07 评估与实施")
    scenarios = [
        "名义200对200",
        "密集交叉",
        "编队分裂与转弯",
        "雷达延迟与虚警",
        "遮挡与外参漂移",
        "高威胁多资源",
        "200对160 / 160对200",
        "中心与二级失效",
        "区域威胁突发",
        "计算与带宽退化",
    ]
    for index, name in enumerate(scenarios):
        column = index % 2
        row = index // 2
        x = 0.62 + column * 4.18
        y = 1.72 + row * 0.88
        add_rect(slide, x, y, 3.72, 0.62, fill=WHITE, line=LINE)
        add_text(slide, f"{index + 1:02d}", x + 0.14, y + 0.15, 0.40, 0.24, size=11, color=TEAL, bold=True, margin=0)
        add_text(slide, name, x + 0.62, y + 0.13, 2.88, 0.28, size=13, color=INK, bold=True, margin=0)
    add_rect(slide, 9.08, 1.72, 3.64, 4.15, fill=TEAL_LIGHT, line=TEAL)
    add_text(slide, "配对原则", 9.38, 2.02, 3.04, 0.35, size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_rich_text(
        slide,
        [
            ("每类场景至少20个未见seed", False, INK),
            ("规则组与候选组使用相同初态", False, INK),
            ("共享传感器、通信和故障随机流", False, INK),
            ("记录模型实际采用，不能只看加载成功", False, INK),
            ("Bootstrap置信区间和失败原因共同判定", False, INK),
        ],
        9.37,
        2.58,
        3.02,
        2.52,
        size=13.5,
        bullet=True,
        gap=10,
    )
    add_label(slide, "安全条件：零真值泄漏、零过期计划、零重复分配", 9.38, 5.32, fill=RED_LIGHT, color=RED, width=3.02)
    add_text(
        slide,
        "候选组必须在不恶化高威胁需求、身份、安全和数值稳定的前提下，至少改善一项预注册任务指标。",
        0.80,
        6.43,
        11.78,
        0.34,
        size=12.5,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 34 Current evidence
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "当前完成度", 34, section="08 结论与下一步")
    add_rect(slide, 0.60, 1.72, 5.95, 4.70, fill=GREEN_LIGHT, line=GREEN)
    add_text(slide, "已形成的证据", 0.92, 2.02, 5.32, 0.38, size=20, color=GREEN, bold=True)
    add_rich_text(
        slide,
        [
            ("200目标、200资源、8侦察节点的10秒三维回合可运行", False, INK),
            ("在线匿名观测12087条，在线真值使用为0", False, INK),
            ("D1-D7均有在线发布，状态有限", False, INK),
            ("D3第8版形成200个唯一资源与200个唯一航迹绑定", False, INK),
            ("计划确认10次，绑定与控制采用1978次", False, INK),
            ("相机命令18928条全部确认执行", False, INK),
        ],
        0.92,
        2.62,
        5.16,
        2.94,
        size=14,
        bullet=True,
        gap=10,
    )
    add_rect(slide, 6.80, 1.72, 5.93, 4.70, fill=RED_LIGHT, line=RED)
    add_text(slide, "尚未闭合", 7.12, 2.02, 5.30, 0.38, size=20, color=RED, bold=True)
    add_rich_text(
        slide,
        [
            ("200规模实时因子达到1", False, INK),
            ("严格身份连续、正式位置误差和完整置信区间", False, INK),
            ("二级节点与完全分布式的稳定全局注册", False, INK),
            ("多随机种子五米物理拦截", False, INK),
            ("学习策略相对规则基线的稳定收益", False, INK),
            ("硬件在环、实飞和真实可见光/红外/雷达联合验证", False, INK),
        ],
        7.12,
        2.62,
        5.16,
        2.94,
        size=14,
        bullet=True,
        gap=10,
    )

    # 35 Risks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "主要风险与处理途径", 35, section="08 结论与下一步")
    risk_rows = [
        ("实时性", "实时因子0.1247", "批量数值、稀疏候选、对象与证据路径优化"),
        ("身份误差放大", "上游混轨传递到分配与导引", "保留协方差、歧义和可用性；错误时保持"),
        ("奖励偏差", "平均收益掩盖高威胁和安全损失", "分解奖励；身份、版本和安全继续用硬约束"),
        ("仿真到实装", "真实误差具有相关性", "真实回放辨识噪声、时钟、外参与运动模糊"),
        ("通信与分区", "200路原始视频不可持续集中上传", "边缘检测与跟踪，只传候选、关键帧和摘要"),
    ]
    for index, (risk, phenomenon, route) in enumerate(risk_rows):
        y = 1.68 + index * 0.95
        color = [RED, ORANGE, PURPLE, BLUE, TEAL][index]
        fill = [RED_LIGHT, ORANGE_LIGHT, PURPLE_LIGHT, BLUE_LIGHT, TEAL_LIGHT][index]
        add_rect(slide, 0.62, y, 12.10, 0.74, fill=WHITE, line=LINE)
        add_rect(slide, 0.62, y, 1.48, 0.74, fill=fill, line=color)
        add_text(slide, risk, 0.72, y + 0.19, 1.28, 0.28, size=13, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, phenomenon, 2.34, y + 0.17, 3.55, 0.30, size=12.5, color=INK, bold=True, margin=0)
        add_text(slide, route, 6.08, y + 0.14, 6.30, 0.36, size=12.2, color=MUTED, margin=0)
    add_text(
        slide,
        "风险处理原则：先保证身份、版本、资源守恒和数值稳定，再讨论速度与学习收益。",
        1.02,
        6.55,
        11.24,
        0.32,
        size=13,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 36 Roadmap & conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "实施路径与阶段结论", 36, section="08 结论与下一步")
    stages = [
        ("阶段1", "确定性基线冻结", "20个未见seed\n严格身份与五米结果", TEAL),
        ("阶段2", "规则性能与容量", "20/50/100/200\n规模曲线和周期预算", BLUE),
        ("阶段3", "单学习模块影子", "D3→D4→D5图网络\n→主动视觉", PURPLE),
        ("阶段4", "有限辅助", "只开放有界修正\n最终合同仍确定性", ORANGE),
        ("阶段5", "AirSim与真实回放", "坐标、相机、通信\n软件/硬件在环", RED),
    ]
    for index, (stage, title, body, color) in enumerate(stages):
        x = 0.50 + index * 2.56
        add_rect(slide, x, 1.70, 2.24, 2.15, fill=WHITE, line=color, width=1.5)
        add_label(slide, stage, x + 0.57, 1.91, fill=color, color=WHITE, width=1.10)
        add_text(slide, title, x + 0.16, 2.49, 1.92, 0.34, size=14.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.18, 2.96, 1.88, 0.62, size=11.5, color=INK, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if index < 4:
            add_arrow(slide, x + 2.25, 2.80, x + 2.48, 2.80, color=MUTED)
    add_rect(slide, 0.72, 4.38, 11.88, 1.38, fill=TEAL_LIGHT, line=TEAL, width=1.5)
    add_text(
        slide,
        "阶段结论",
        1.02,
        4.72,
        1.55,
        0.34,
        size=18,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "200对200确定性模块栈已完成单回合规模化接口闭环。\n"
        "近期重点是多随机种子基线、严格身份谱系和端到端性能。\n"
        "学习增强逐项进入影子验证，不同时改变在线计划。",
        2.72,
        4.63,
        9.43,
        0.70,
        size=16,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_label(slide, "当前重点：阶段1—2", 5.45, 6.22, fill=ORANGE, color=WHITE, width=2.42)

    assert len(prs.slides) == 36
    return prs


def main() -> None:
    presentation = build_deck()
    presentation.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
