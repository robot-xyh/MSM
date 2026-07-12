#!/usr/bin/env python3
"""Generate the Chinese leadership briefing deck for the MSM C-UAS project."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentations"
ASSET_DIR = OUT_DIR / "assets"
OUTPUT = OUT_DIR / "C_UAS_MULTI_UAV_INTERCEPTION_LEADERSHIP_BRIEFING_CN.pptx"

W = 13.333
H = 7.5

FONT = "Noto Sans CJK SC"
FONT_MONO = "DejaVu Sans Mono"

WHITE = "FFFFFF"
BG = "FFFFFF"
INK = "1C2733"
MUTED = "596673"
LINE = "C8D0D8"
RED = "A83B3B"
TEAL = "28777C"
BLUE = "315B82"
GREEN = "477358"
AMBER = "A97422"
PURPLE = "62556F"
PALE_RED = "F4EAEA"
PALE_TEAL = "EAF2F2"
PALE_BLUE = "EAF0F5"
PALE_GREEN = "EBF1ED"
PALE_AMBER = "F5F0E6"
PALE_PURPLE = "EFECF2"
DARK = "263746"


ASSETS = {
    "d1_tracks": ROOT / "research_modules/d1_sensor_fusion/reports/tracks_xy.png",
    "d1_latency": ROOT / "research_modules/d1_sensor_fusion/reports/rmse_latency_ablation.png",
    "d2_result": ROOT / "research_modules/d2_data_association/docs/association_idsw_rmse.png",
    "d3_result": ROOT / "research_modules/d3_assignment_planner/results/cost_reassignment.png",
    "d3_weights": ROOT / "research_modules/d3_assignment_planner/results/weight_sensitivity.png",
    "d4_curve": ROOT / "research_modules/d4_distributed_fallback/reports/failover_packet_loss_curve.png",
    "d5_timeline": ROOT / "research_modules/d5_terminal_association/docs/terminal_decision_timeline.png",
    "cv_accept": ROOT / "research_modules/airsim_runtime/outputs/p1_p2_validation_20260711/p1_cv_acceptance_by_seed.png",
    "sf_diag": ROOT / "research_modules/airsim_runtime/outputs/p1_p2_validation_20260711/p1_simpleflight_diagnostic.png",
    "p2_guidance": ROOT / "research_modules/airsim_runtime/outputs/p1_p2_validation_20260711/p2_guidance_miss_distance.png",
    "global_gif": ROOT / "research_modules/integrated_simulation/outputs/global_process_2d.gif",
    "air_3d": ROOT / "research_modules/airsim_runtime/outputs/p1_gap_closure_2v2_multiseed_20260710_seed002/episode_006_full_flow/airsim_3d_intercept_trajectories.png",
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color: str, width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True, line_width=1.0):
    # Engineering-report visual language: square boxes, thin rules, no decorative pills.
    kind = MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, line_width)
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    font: str = FONT,
    margin: float = 0.03,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, segments, x, y, w, h, size=18, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for seg in segments:
        run = p.add_run()
        run.text = seg.get("text", "")
        run.font.name = seg.get("font", FONT)
        run.font.size = Pt(seg.get("size", size))
        run.font.bold = seg.get("bold", False)
        run.font.color.rgb = rgb(seg.get("color", INK))
    return box


def add_bullets(slide, items: Iterable[str], x, y, w, h, size=16, color=INK, leading=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(size * 0.45)
        p.line_spacing = leading
    return box


def add_title(slide, title: str, subtitle: str | None = None, section: str | None = None):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(0.12))
    set_fill(header, DARK)
    header.line.fill.background()
    add_text(slide, "MSM C-UAS · 阶段技术汇报", 0.55, 0.27, 2.8, 0.28, 9.2, BLUE, True)
    add_text(slide, title, 0.55, 0.56, 12.1, 0.52, 25, INK, True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.07, 11.9, 0.36, 11.5, MUTED)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.43), Inches(0.65), Inches(0.05))
    set_fill(accent, RED)
    accent.line.fill.background()


def add_footer(slide, number: int, label: str = "MSM C-UAS · 阶段技术汇报"):
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.08), Inches(12.15), Inches(0.01))
    set_fill(rule, LINE)
    rule.line.fill.background()
    add_text(slide, label, 0.55, 7.18, 4.2, 0.18, 8.5, MUTED)
    add_text(slide, f"{number:02d}", 12.1, 7.16, 0.62, 0.2, 8.5, MUTED, True, PP_ALIGN.RIGHT)


def new_slide(prs, title=None, subtitle=None, section=None, bg=BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(bg)
    if title:
        add_title(slide, title, subtitle, section)
    return slide


def add_badge(slide, text, x, y, w, color=TEAL, fill=PALE_TEAL, size=10.5):
    s = add_rect(slide, x, y, w, 0.34, fill, fill, True, 0)
    add_text(slide, text, x + 0.05, y + 0.03, w - 0.1, 0.24, size, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return s


def add_metric(slide, value, label, x, y, w=1.65, color=TEAL, note=None):
    add_rect(slide, x, y, w, 1.15, WHITE, LINE)
    add_text(slide, value, x + 0.08, y + 0.12, w - 0.16, 0.48, 25, color, True, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.08, y + 0.64, w - 0.16, 0.28, 10.5, MUTED, False, PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, x + 0.08, y + 0.91, w - 0.16, 0.16, 7.5, MUTED, False, PP_ALIGN.CENTER)


def add_card(slide, title, body, x, y, w, h, accent=TEAL, fill=WHITE, title_size=15, body_size=11.5):
    add_rect(slide, x, y, w, h, fill, LINE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    set_fill(bar, accent)
    bar.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.32, 0.35, title_size, INK, True)
    add_text(slide, body, x + 0.18, y + 0.57, w - 0.32, h - 0.7, body_size, MUTED)


def add_node(slide, title, subtitle, x, y, w, h, color=TEAL, fill=WHITE, title_size=14):
    add_rect(slide, x, y, w, h, fill, color, True, 1.6)
    add_text(slide, title, x + 0.08, y + 0.11, w - 0.16, 0.32, title_size, color, True, PP_ALIGN.CENTER)
    add_text(slide, subtitle, x + 0.08, y + 0.46, w - 0.16, h - 0.53, 9.5, MUTED, False, PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, width=1.5, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    set_line(line, color, width)
    line.line.end_arrowhead = True
    if dashed:
        line.line.dash_style = 2
    return line


def add_picture_contain(slide, path: Path, x, y, w, h, border=True, bg=WHITE):
    if not path.exists():
        add_rect(slide, x, y, w, h, PALE_RED, RED)
        add_text(slide, f"素材缺失\n{path.name}", x + 0.2, y + h / 2 - 0.35, w - 0.4, 0.7, 12, RED, True, PP_ALIGN.CENTER)
        return None
    if border:
        add_rect(slide, x, y, w, h, bg, LINE)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min((w - 0.08) / iw, (h - 0.08) / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def add_status_dot(slide, x, y, color, label, value=None):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.18), Inches(0.18))
    set_fill(dot, color)
    dot.line.fill.background()
    add_text(slide, label, x + 0.27, y - 0.03, 2.55, 0.25, 11, INK, value is not None)
    if value is not None:
        add_text(slide, value, x + 2.15, y - 0.03, 1.2, 0.25, 11, color, True, PP_ALIGN.RIGHT)


def add_phase_band(slide, x, y, w, title, subtitle, color, fill):
    add_rect(slide, x, y, w, 0.88, fill, color, True, 1.2)
    add_text(slide, title, x + 0.12, y + 0.12, w - 0.24, 0.28, 13, color, True, PP_ALIGN.CENTER)
    add_text(slide, subtitle, x + 0.12, y + 0.45, w - 0.24, 0.25, 9.5, MUTED, False, PP_ALIGN.CENTER)


def add_module_banner(slide, code, name, value, x=0.65, y=1.65, color=TEAL, fill=PALE_TEAL):
    add_rect(slide, x, y, 2.1, 0.86, fill, color)
    add_text(slide, code, x + 0.14, y + 0.12, 0.55, 0.45, 22, color, True, PP_ALIGN.CENTER)
    add_text(slide, name, x + 0.72, y + 0.1, 1.25, 0.26, 12.5, INK, True)
    add_text(slide, value, x + 0.72, y + 0.42, 1.25, 0.23, 9, MUTED)


def create_deck() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 01 Cover - same visual system as the internal pages.
    slide = new_slide(prs, bg=WHITE)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(0.16))
    set_fill(header, DARK)
    header.line.fill.background()
    add_text(slide, "MSM 反无人机多无人机拦截科研仿真项目", 0.7, 0.55, 6.2, 0.35, 12, BLUE, True)
    add_text(slide, "多无人机反无人机拦截体系", 0.7, 1.15, 6.0, 0.72, 30, INK, True)
    add_text(slide, "阶段技术方案与仿真验证汇报", 0.72, 2.02, 5.2, 0.48, 19, MUTED)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(2.66), Inches(0.85), Inches(0.06))
    set_fill(accent, RED)
    accent.line.fill.background()
    add_rect(slide, 0.72, 3.2, 5.3, 2.15, "F6F7F8", LINE, False)
    rows = [
        ("汇报范围", "总体架构、D1-D7 模块、AirSim 验证、下一阶段计划"),
        ("当前阶段", "P1 合同层闭合，物理拦截闭环专项验证中"),
        ("验证环境", "质点模型 + AirSim ComputerVision / SimpleFlight"),
        ("汇报日期", "2026 年 7 月"),
    ]
    for idx, (label, value) in enumerate(rows):
        y = 3.42 + idx * 0.44
        add_text(slide, label, 0.95, y, 1.0, 0.24, 10.5, MUTED, True)
        add_text(slide, value, 2.02, y, 3.7, 0.28, 10.5, INK)
        if idx < len(rows) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y + 0.32), Inches(4.75), Inches(0.01))
            set_fill(line, LINE)
            line.line.fill.background()
    add_picture_contain(slide, ASSETS["air_3d"], 6.55, 0.72, 6.0, 5.9, True, WHITE)
    add_text(slide, "图：AirSim 受控拦截三维轨迹（历史验证样例）", 7.0, 6.65, 5.1, 0.24, 9, MUTED, False, PP_ALIGN.CENTER)
    add_text(slide, "说明：本项目处于科研仿真阶段，不涉及实装处置与自动授权。", 0.75, 6.65, 5.4, 0.24, 9.5, MUTED)
    add_footer(slide, 1, "MSM C-UAS · 阶段技术汇报")

    # 02 Project summary in engineering-report format.
    slide = new_slide(prs, "项目概况与阶段结论", "本页数据均来自当前代码、模块测试和 AirSim 运行产物", "PROJECT STATUS")
    add_text(slide, "1. 项目概况", 0.7, 1.72, 2.1, 0.34, 15, INK, True)
    summary_rows = [
        ("建设目标", "建立多目标、多资源、可降级的反无人机拦截科研仿真体系"),
        ("系统组成", "中心 C2、高空侦察二级节点、拦截无人机集群"),
        ("技术主线", "多源融合、目标关联、动态分配、降级协同、视觉配准、比例导引、统一评估"),
        ("运行方式", "正常中心化；中心失效后由二级节点接管；二级失效后完全分布式"),
        ("验证方式", "质点模型验证算法闭环；AirSim 验证接口、视觉、状态机与 SimpleFlight 控制"),
    ]
    add_rect(slide, 0.7, 2.15, 7.0, 3.35, WHITE, LINE, False)
    for idx, (label, value) in enumerate(summary_rows):
        y = 2.16 + idx * 0.66
        fill = "F2F4F6" if idx % 2 == 0 else WHITE
        add_rect(slide, 0.71, y, 1.35, 0.64, fill, LINE, False, 0.5)
        add_rect(slide, 2.06, y, 5.63, 0.64, fill, LINE, False, 0.5)
        add_text(slide, label, 0.82, y + 0.18, 1.12, 0.24, 10.5, BLUE, True, PP_ALIGN.CENTER)
        add_text(slide, value, 2.22, y + 0.15, 5.25, 0.32, 10.5, INK)
    add_text(slide, "2. 阶段状态", 8.05, 1.72, 2.1, 0.34, 15, INK, True)
    status_rows = [
        ("P1 合同层", "完成", GREEN),
        ("CV 双 primary", "8/10", GREEN),
        ("二级/peer 联盟", "ACK 3/3", TEAL),
        ("ID Switch / 错误重复锁", "0", BLUE),
        ("SimpleFlight 物理命中", "0/30", RED),
        ("末端检测超时", "24/30", RED),
        ("P2 外部算法", "隔离对照", PURPLE),
    ]
    add_rect(slide, 8.05, 2.15, 4.55, 3.35, WHITE, LINE, False)
    for idx, (label, value, color) in enumerate(status_rows):
        y = 2.18 + idx * 0.46
        add_text(slide, label, 8.28, y + 0.05, 2.6, 0.24, 10.2, INK)
        add_text(slide, value, 10.92, y + 0.05, 1.35, 0.24, 10.5, color, True, PP_ALIGN.RIGHT)
        if idx < len(status_rows) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.25), Inches(y + 0.34), Inches(4.05), Inches(0.01))
            set_fill(line, LINE)
            line.line.fill.background()
    add_rect(slide, 0.7, 5.85, 11.9, 0.72, "F2F4F6", LINE, False)
    add_text(slide, "下一阶段主任务", 0.95, 6.05, 1.75, 0.28, 12, BLUE, True)
    add_text(slide, "完成末端检测持续性、D5 锁定、D7 切换和物理接近闭环的分层标定。", 2.75, 6.02, 9.35, 0.32, 13, INK, True)
    add_footer(slide, 2)

    # 03 Scenario
    slide = new_slide(prs, "任务场景与主要工程问题", "目标多、链路异步、末端视场复杂，同时需要考虑中心节点失效", "SCENARIO")
    add_text(slide, "威胁侧", 0.75, 1.75, 1.1, 0.3, 13, RED, True)
    for i, (x, y) in enumerate([(1.0, 2.4), (1.55, 3.25), (0.92, 4.2), (1.75, 4.9)]):
        shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(x), Inches(y), Inches(0.62), Inches(0.46))
        set_fill(shape, RED)
        set_line(shape, RED)
        add_text(slide, f"T{i+1}", x, y + 0.09, 0.62, 0.2, 9, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, "多目标交叉 / 机动 / 遮挡", 0.55, 5.6, 2.6, 0.35, 11, MUTED, False, PP_ALIGN.CENTER)
    add_node(slide, "中心节点", "全局航迹\n统一分配\n授权与审计", 5.15, 2.0, 2.1, 1.35, BLUE, PALE_BLUE)
    add_node(slide, "高空侦察节点", "全局光电视野\n二级接管\n区域视频/数据", 5.15, 4.35, 2.1, 1.35, TEAL, PALE_TEAL)
    add_text(slide, "拦截资源", 9.25, 1.75, 1.2, 0.3, 13, GREEN, True)
    for i, (x, y) in enumerate([(9.05, 2.35), (10.1, 2.95), (9.0, 4.1), (10.15, 4.8)]):
        add_node(slide, f"INT-{i+1}", "雷达中段 PN\n相机末端 PNG", x, y, 1.45, 0.78, GREEN, PALE_GREEN, 10.5)
    add_arrow(slide, 2.45, 3.6, 5.05, 2.7, RED, 1.5)
    add_arrow(slide, 7.3, 2.7, 8.95, 3.2, BLUE, 1.5)
    add_arrow(slide, 7.3, 4.9, 8.95, 4.55, TEAL, 1.5)
    add_arrow(slide, 6.2, 3.38, 6.2, 4.25, MUTED, 1.2, True)
    add_card(slide, "主要约束", "在不确定、异步和局部视野条件下，保持目标身份、任务版本和控制授权一致。", 3.1, 5.9, 7.1, 0.72, AMBER, PALE_AMBER, 13, 11)
    add_footer(slide, 3)

    # 04 Three-tier architecture
    slide = new_slide(prs, "三层节点体系", "中心最优、二级连续、分布式保底", "SYSTEM CONCEPT")
    layers = [
        ("一级：中心 C2", "多传感器融合、全局航迹、动态分配、授权与评估", BLUE, PALE_BLUE, 1.75),
        ("二级：机动高空侦察节点", "高性能光电云台、区域全局视野、二级计划 owner、视频/数据中继", TEAL, PALE_TEAL, 3.18),
        ("三级：拦截无人机集群", "局部视觉、多机通信、完全分布式联盟、位置 PN 与视觉 PNG", GREEN, PALE_GREEN, 4.62),
    ]
    for title, body, color, fill, y in layers:
        add_rect(slide, 0.9, y, 11.55, 1.05, fill, color, True, 1.5)
        add_text(slide, title, 1.15, y + 0.18, 3.25, 0.36, 17, color, True)
        add_text(slide, body, 4.2, y + 0.18, 7.8, 0.45, 13, INK)
    add_arrow(slide, 6.65, 2.83, 6.65, 3.1, MUTED, 1.6)
    add_arrow(slide, 6.65, 4.26, 6.65, 4.54, MUTED, 1.6)
    add_badge(slide, "正常：中心主控", 1.05, 6.15, 2.4, BLUE, PALE_BLUE)
    add_badge(slide, "中心失效：二级接管", 3.8, 6.15, 2.8, TEAL, PALE_TEAL)
    add_badge(slide, "二级失效：完全分布式", 6.95, 6.15, 3.0, GREEN, PALE_GREEN)
    add_badge(slide, "全程版本化 + fail-closed", 10.2, 6.15, 2.2, RED, PALE_RED, 9.5)
    add_footer(slide, 4)

    # 05 Overall architecture
    slide = new_slide(prs, "D1-D7 端到端技术架构", "每个模块只做一件事，统一通过带时间戳、协方差和版本的数据合同协同", "ARCHITECTURE")
    modules = [
        ("D1", "多源融合", "GlobalTrack", TEAL, PALE_TEAL),
        ("D2", "数据关联", "稳定全局ID", BLUE, PALE_BLUE),
        ("D3", "资源分配", "AssignmentPlan", AMBER, PALE_AMBER),
        ("D4", "降级协同", "C2/联盟状态", PURPLE, PALE_PURPLE),
        ("D5", "视觉配准", "TerminalAssociation", RED, PALE_RED),
        ("D7", "比例导引", "PN → PNG", GREEN, PALE_GREEN),
    ]
    xs = [0.65, 2.7, 4.75, 6.8, 8.85, 10.9]
    for i in range(len(xs) - 1):
        add_arrow(slide, xs[i] + 1.65, 3.05, xs[i + 1] - 0.08, 3.05, MUTED, 1.6)
    for (code, name, output, color, fill), x in zip(modules, xs):
        add_node(slide, f"{code}  {name}", output, x, 2.55, 1.7, 1.05, color, fill, 12.5)
    add_node(slide, "D6 系统评估", "消费全链路日志\n不参与控制", 5.35, 4.65, 2.65, 1.0, INK, WHITE, 13)
    for x in xs:
        add_arrow(slide, x + 0.85, 3.7, 6.65, 4.58, LINE, 0.9)
    add_rect(slide, 0.85, 1.67, 11.65, 0.45, WHITE, LINE)
    add_text(slide, "雷达 / 声学 / 光电 / LiDAR（可选）", 1.1, 1.78, 3.1, 0.25, 11, MUTED, True, PP_ALIGN.CENTER)
    add_text(slide, "统一 NED / 双时间戳 / 协方差 / global_track_id / 计划版本", 4.25, 1.78, 7.9, 0.25, 11, INK, True, PP_ALIGN.CENTER)
    add_card(slide, "接口约束", "D5/D7 不改写 global_track_id；拒绝过期计划；证据不足时保持 hold，不允许自行换绑。", 2.35, 6.05, 8.65, 0.7, RED, PALE_RED, 13, 11)
    add_footer(slide, 5)

    # 06 Bus and communications
    slide = new_slide(prs, "通信拓扑与统一数据总线", "视频用于发现与复核，结构化摘要用于决策；控制权限始终受版本和身份约束", "DATA & COMMUNICATION")
    add_node(slide, "中心 C2", "GlobalTrack\nAssignmentPlan\n授权/审计", 5.25, 1.72, 2.25, 1.05, BLUE, PALE_BLUE)
    add_node(slide, "高空侦察 1", "区域视频\nReconImageCue\n二级计划", 2.0, 3.22, 2.3, 1.1, TEAL, PALE_TEAL)
    add_node(slide, "高空侦察 2", "区域视频\nTrackSummary\n通信中继", 9.0, 3.22, 2.3, 1.1, TEAL, PALE_TEAL)
    for i, x in enumerate([0.75, 3.4, 6.75, 9.4]):
        add_node(slide, f"拦截机 {i+1}", "LocalVisualTrack\nResourceSummary", x, 5.3, 2.05, 0.9, GREEN, PALE_GREEN, 11.5)
    add_arrow(slide, 5.2, 2.45, 4.35, 3.22, BLUE, 1.2)
    add_arrow(slide, 7.55, 2.45, 8.95, 3.22, BLUE, 1.2)
    for x in [1.75, 4.4, 7.75, 10.4]:
        source_x = 3.15 if x < 6 else 10.15
        add_arrow(slide, source_x, 4.35, x, 5.25, TEAL, 1.0)
    add_arrow(slide, 2.85, 5.75, 3.35, 5.75, GREEN, 1.0, True)
    add_arrow(slide, 5.5, 5.75, 6.7, 5.75, GREEN, 1.0, True)
    add_arrow(slide, 8.85, 5.75, 9.35, 5.75, GREEN, 1.0, True)
    add_badge(slide, "TrackSummary", 0.85, 6.55, 1.65, BLUE, PALE_BLUE)
    add_badge(slide, "ResourceSummary", 2.7, 6.55, 1.9, GREEN, PALE_GREEN)
    add_badge(slide, "IdentityClaim", 4.8, 6.55, 1.7, RED, PALE_RED)
    add_badge(slide, "CoalitionCommit", 6.7, 6.55, 2.0, PURPLE, PALE_PURPLE)
    add_badge(slide, "EventRecord", 8.9, 6.55, 1.55, AMBER, PALE_AMBER)
    add_badge(slide, "视频/图像 cue", 10.65, 6.55, 1.7, TEAL, PALE_TEAL)
    add_footer(slide, 6)

    # 07 Degradation workflow
    slide = new_slide(prs, "正常运行、主动降级与被动降级", "主动降级由不确定性和末端不一致触发；被动降级由节点失效触发", "OPERATING MODES")
    add_phase_band(slide, 0.75, 1.78, 2.35, "中心正常", "中心计划持续有效", BLUE, PALE_BLUE)
    add_phase_band(slide, 3.45, 1.78, 2.35, "主动降级判断", "D1/D2不确定度 + D5不一致", AMBER, PALE_AMBER)
    add_phase_band(slide, 6.15, 1.78, 2.35, "二级节点接管", "高空侦察节点发布新版本", TEAL, PALE_TEAL)
    add_phase_band(slide, 8.85, 1.78, 2.35, "完全分布式", "peer联盟 ACK/commit", GREEN, PALE_GREEN)
    add_arrow(slide, 3.12, 2.22, 3.42, 2.22, MUTED, 1.6)
    add_arrow(slide, 5.82, 2.22, 6.12, 2.22, MUTED, 1.6)
    add_arrow(slide, 8.52, 2.22, 8.82, 2.22, MUTED, 1.6)
    add_card(slide, "继续执行", "中心/二级分配与 D5 末端关联一致，计划版本新鲜，联盟证据有效。", 0.85, 3.35, 3.5, 1.2, GREEN, PALE_GREEN)
    add_card(slide, "请求重规划", "中心仍在线，但跟踪不确定性、计划延时或视觉证据表明原计划失效。", 4.9, 3.35, 3.5, 1.2, AMBER, PALE_AMBER)
    add_card(slide, "保守阻断", "缺 ACK、旧 epoch、过期 lease、身份冲突或网络分区时 fail-closed。", 8.95, 3.35, 3.5, 1.2, RED, PALE_RED)
    add_rect(slide, 1.2, 5.15, 10.9, 1.0, WHITE, LINE)
    add_text(slide, "切换规则", 1.45, 5.42, 1.2, 0.3, 13, PURPLE, True)
    add_text(slide, "结果一致时继续执行；结果不一致时进入仲裁；证据不足时不授权；中心恢复后先完成双轨校验。", 2.55, 5.35, 8.95, 0.42, 15, INK, True, PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 08 D1 intro
    slide = new_slide(prs, "D1 多传感器融合与目标配准", "把异步、异构、不同参考系的观测统一成可信 GlobalTrack", "MODULE D1")
    add_module_banner(slide, "D1", "多源融合", "观测 → GlobalTrack", color=TEAL, fill=PALE_TEAL)
    add_node(slide, "雷达", "三维位置/速度\n距离相关协方差", 0.8, 3.0, 2.0, 0.95, BLUE, PALE_BLUE)
    add_node(slide, "声学", "粗方位 + 声纹\n低空辅助", 0.8, 4.25, 2.0, 0.95, AMBER, PALE_AMBER)
    add_node(slide, "光电", "像素框 + 相机参数\n局部确认", 0.8, 5.5, 2.0, 0.95, RED, PALE_RED)
    add_arrow(slide, 2.85, 3.48, 4.0, 4.28, MUTED, 1.2)
    add_arrow(slide, 2.85, 4.72, 4.0, 4.5, MUTED, 1.2)
    add_arrow(slide, 2.85, 5.98, 4.0, 4.72, MUTED, 1.2)
    add_node(slide, "时空对齐", "measurement / arrival 双时间戳\nsensor frame → NED", 4.05, 3.75, 2.45, 1.45, TEAL, PALE_TEAL)
    add_arrow(slide, 6.55, 4.48, 7.35, 4.48, MUTED, 1.4)
    add_node(slide, "轻量 EKF / OOSM", "延迟补偿\n协方差传播\nsource lineage 去重", 7.4, 3.75, 2.4, 1.45, PURPLE, PALE_PURPLE)
    add_arrow(slide, 9.85, 4.48, 10.5, 4.48, MUTED, 1.4)
    add_node(slide, "GlobalTrack", "position / velocity\ncovariance\ntrack quality", 10.55, 3.75, 2.0, 1.45, GREEN, PALE_GREEN)
    add_badge(slide, "coarse", 4.0, 5.75, 1.25, AMBER, PALE_AMBER)
    add_badge(slide, "stable", 5.5, 5.75, 1.25, TEAL, PALE_TEAL)
    add_badge(slide, "handover", 7.0, 5.75, 1.45, GREEN, PALE_GREEN)
    add_card(slide, "价值", "后续关联、分配、视觉门控与安全决策都使用“均值 + 不确定性”，而不是只追一个点。", 8.7, 5.7, 3.75, 0.8, TEAL, PALE_TEAL, 12.5, 10.5)
    add_footer(slide, 8)

    # 09 D1 results
    slide = new_slide(prs, "D1 验证结果", "延迟补偿与融合结果可量化，当前主线保持轻量可解释", "MODULE D1")
    add_picture_contain(slide, ASSETS["d1_tracks"], 0.65, 1.72, 6.1, 4.72)
    add_picture_contain(slide, ASSETS["d1_latency"], 6.95, 1.72, 5.72, 3.35)
    add_card(slide, "已完成", "双时间戳、NED、协方差、延迟/OOSM、距离相关噪声、航迹分级、governed replay。", 7.0, 5.27, 2.7, 1.15, GREEN, PALE_GREEN, 13, 10.5)
    add_card(slide, "仍需标定", "真实长 replay、复杂遮挡、多节点公共信息相关性和长期阈值治理。", 9.9, 5.27, 2.7, 1.15, AMBER, PALE_AMBER, 13, 10.5)
    add_text(slide, "图：D1 质点融合航迹与延迟补偿消融", 0.8, 6.55, 6.0, 0.25, 9, MUTED, False, PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # 10 D2 intro
    slide = new_slide(prs, "D2 多目标跟踪与数据关联", "目标交叉时保持身份连续，避免重复分配和错误末端锁定", "MODULE D2")
    add_module_banner(slide, "D2", "航迹关联", "稳定 global_track_id", color=BLUE, fill=PALE_BLUE)
    add_text(slide, "预测航迹", 0.9, 3.0, 1.2, 0.28, 12, BLUE, True)
    add_text(slide, "新观测", 0.9, 5.05, 1.2, 0.28, 12, RED, True)
    for i, y in enumerate([3.45, 4.0, 4.55]):
        add_node(slide, f"Track {i+1}", "状态 + 协方差", 1.0, y, 1.65, 0.42, BLUE, PALE_BLUE, 9.5)
    for i, y in enumerate([5.48, 5.95, 6.42]):
        add_node(slide, f"Detection {i+1}", "measurement", 1.0, y, 1.65, 0.36, RED, PALE_RED, 8.8)
    add_node(slide, "马氏门限", "排除不可能候选", 3.55, 3.62, 1.8, 0.92, TEAL, PALE_TEAL)
    add_node(slide, "代价矩阵", "运动 + 类别 + 时间", 5.85, 3.62, 1.9, 0.92, AMBER, PALE_AMBER)
    add_node(slide, "GNN / Hungarian", "默认硬关联主线", 8.3, 3.62, 2.05, 0.92, BLUE, PALE_BLUE)
    add_node(slide, "航迹生命周期", "tentative → confirmed\nengageable → lost", 10.85, 3.5, 1.75, 1.15, GREEN, PALE_GREEN, 11)
    add_arrow(slide, 2.7, 4.0, 3.48, 4.0, MUTED, 1.2)
    add_arrow(slide, 2.7, 5.8, 3.48, 4.28, MUTED, 1.2)
    add_arrow(slide, 5.4, 4.08, 5.8, 4.08, MUTED, 1.2)
    add_arrow(slide, 7.8, 4.08, 8.25, 4.08, MUTED, 1.2)
    add_arrow(slide, 10.4, 4.08, 10.8, 4.08, MUTED, 1.2)
    add_badge(slide, "默认：GNN/Hungarian", 3.6, 5.25, 2.45, BLUE, PALE_BLUE)
    add_badge(slide, "升级：JPDA", 6.3, 5.25, 1.65, TEAL, PALE_TEAL)
    add_badge(slide, "研究对照：MHT", 8.2, 5.25, 1.85, PURPLE, PALE_PURPLE)
    add_card(slide, "强制指标", "id_switch_count、track_continuity、duplicate_assignment_count 全程显式记录。", 4.0, 6.05, 6.2, 0.65, RED, PALE_RED, 12.5, 10.5)
    add_footer(slide, 10)

    # 11 D2 result
    slide = new_slide(prs, "D2 验证结果", "身份连续性已经成为系统级硬指标，而不是事后观察项", "MODULE D2")
    add_picture_contain(slide, ASSETS["d2_result"], 0.65, 1.7, 7.45, 4.95)
    add_metric(slide, "0", "CV 10-seed ID Switch", 8.55, 1.95, 1.85, BLUE)
    add_metric(slide, "1.0", "P2 replay continuity", 10.6, 1.95, 1.85, GREEN)
    add_card(slide, "主线路径", "NumPy/SciPy GNN/Hungarian 端到端 replay 可运行，保持轻量和可解释。", 8.55, 3.55, 3.9, 1.1, BLUE, PALE_BLUE, 13, 10.5)
    add_card(slide, "P2 边界", "Stone Soup / FilterPy 目前只是对象 adapter smoke，不是完整 JPDA/MHT tracker。", 8.55, 4.9, 3.9, 1.1, PURPLE, PALE_PURPLE, 13, 10.5)
    add_text(slide, "下一步：真实长时 crossing / 遮挡 / 漏检 replay 下校准风险阈值", 8.6, 6.27, 3.8, 0.4, 11, AMBER, True, PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12 D3 intro
    slide = new_slide(prs, "D3 集中式资源—目标分配", "从一次性 N 对 N，升级为动态 M 对 N、带迟滞和角色的滚动计划", "MODULE D3")
    add_module_banner(slide, "D3", "动态分配", "GlobalTrack → AssignmentPlan", color=AMBER, fill=PALE_AMBER)
    add_text(slide, "代价组成", 3.25, 1.78, 1.2, 0.3, 13, AMBER, True)
    costs = [
        ("拦截窗口", BLUE, PALE_BLUE), ("航迹不确定性", TEAL, PALE_TEAL),
        ("目标威胁度", RED, PALE_RED), ("资源状态", GREEN, PALE_GREEN),
        ("视场确认难度", PURPLE, PALE_PURPLE), ("资源冲突风险", AMBER, PALE_AMBER),
    ]
    for idx, (name, color, fill) in enumerate(costs):
        x = 3.2 + (idx % 3) * 2.15
        y = 2.25 + (idx // 3) * 0.82
        add_badge(slide, name, x, y, 1.85, color, fill, 9.5)
    add_node(slide, "Demand-slot Hungarian", "按目标需求复制任务槽\n支持 M ≠ N", 4.65, 4.05, 2.65, 1.15, AMBER, PALE_AMBER)
    add_arrow(slide, 7.4, 4.62, 8.35, 4.62, MUTED, 1.4)
    add_node(slide, "版本化计划", "plan id / version / owner\n拒绝 stale plan", 8.4, 4.05, 2.35, 1.15, BLUE, PALE_BLUE)
    add_card(slide, "迟滞策略", "新计划收益必须超过阈值，并满足最小驻留时间，避免目标轻微变化造成频繁换绑。", 0.85, 3.25, 2.6, 1.55, TEAL, PALE_TEAL, 13, 10.5)
    add_card(slide, "角色保持", "primary 健康时保持原角色，仅在 reserve 软失败时局部重解，减少联盟抖动。", 0.85, 5.15, 2.6, 1.35, GREEN, PALE_GREEN, 13, 10.5)
    add_card(slide, "输出", "主资源、备份资源、波次、时间窗、联盟版本和人工授权状态。", 8.25, 5.65, 3.75, 0.75, RED, PALE_RED, 12.5, 10.5)
    add_footer(slide, 12)

    # 13 M-to-N coalition
    slide = new_slide(prs, "高威胁目标的 M 对 N 联盟", "默认研究方案：2 个 primary 同批接近 + 1 个 reserve 待命", "COOPERATIVE INTERCEPTION")
    add_node(slide, "高威胁目标 T001", "required resource count = 3", 5.25, 1.75, 2.8, 0.9, RED, PALE_RED, 14)
    add_node(slide, "Primary A", "第一波\n独立 PN/PNG 状态", 1.1, 3.35, 2.25, 1.0, GREEN, PALE_GREEN)
    add_node(slide, "Primary B", "第一波\n独立 PN/PNG 状态", 4.05, 3.35, 2.25, 1.0, GREEN, PALE_GREEN)
    add_node(slide, "Reserve C", "第二波/补位\n默认 standby", 7.0, 3.35, 2.25, 1.0, AMBER, PALE_AMBER)
    add_node(slide, "普通目标 T002", "单 primary 独立拦截", 10.0, 3.35, 2.15, 1.0, BLUE, PALE_BLUE)
    for x in [2.2, 5.15, 8.1]:
        add_arrow(slide, x, 3.28, 6.65, 2.7, MUTED, 1.25)
    add_card(slide, "原子激活", "全部必要成员 ACK 且 lease 有效，联盟才进入 committed / executing。", 0.9, 5.18, 3.65, 1.15, PURPLE, PALE_PURPLE, 13, 10.5)
    add_card(slide, "协同定位", "多个无人机视角可提供 bearing、bbox 历史和几何质量，但不得重复计算公共信息。", 4.83, 5.18, 3.65, 1.15, TEAL, PALE_TEAL, 13, 10.5)
    add_card(slide, "安全边界", "reserve 未激活不得切换；合法协同多锁不算错误 duplicate；成员退出需重构。", 8.75, 5.18, 3.65, 1.15, RED, PALE_RED, 13, 10.5)
    add_footer(slide, 13)

    # 14 D3 result
    slide = new_slide(prs, "D3 验证结果", "迟滞降低重分配抖动，M 对 N 合同已进入统一运行时", "MODULE D3")
    add_picture_contain(slide, ASSETS["d3_result"], 0.65, 1.72, 7.05, 4.82)
    add_metric(slide, "1.0", "需求满足率（合同场景）", 8.05, 1.95, 2.0, GREEN)
    add_metric(slide, "0", "错误重复分配", 10.3, 1.95, 2.0, BLUE)
    add_card(slide, "P1 已完成", "demand-slot、增量规划、role-aware primary、current-binding 与版本拒绝。", 8.05, 3.55, 4.25, 1.05, GREEN, PALE_GREEN, 13, 10.5)
    add_card(slide, "P1 仍开放", "真实动态 3v5/5v3、目标新增/资源失效、反馈权重和长期迟滞标定。", 8.05, 4.82, 4.25, 1.05, AMBER, PALE_AMBER, 13, 10.5)
    add_badge(slide, "默认：SciPy Hungarian", 8.05, 6.12, 2.25, AMBER, PALE_AMBER)
    add_badge(slide, "P2：OR-Tools 仅接口", 10.55, 6.12, 1.95, PURPLE, PALE_PURPLE, 9.2)
    add_footer(slide, 14)

    # 15 D4 intro
    slide = new_slide(prs, "D4 分布式协同与降级接管", "不是简单“中心掉线后拍卖”，而是有健康判定、接管优先级和原子联盟", "MODULE D4")
    add_module_banner(slide, "D4", "降级协同", "中心 → 二级 → peer", color=PURPLE, fill=PALE_PURPLE)
    states = [("normal", BLUE, PALE_BLUE), ("degraded", AMBER, PALE_AMBER), ("suspect", RED, PALE_RED), ("failed", DARK, "E9ECEF")]
    xs = [3.15, 5.3, 7.45, 9.6]
    for i in range(3):
        add_arrow(slide, xs[i] + 1.55, 2.55, xs[i+1] - 0.08, 2.55, MUTED, 1.2)
    for (name, color, fill), x in zip(states, xs):
        add_node(slide, name, "heartbeat / track / plan timeout", x, 2.1, 1.6, 0.9, color, fill, 11)
    add_text(slide, "C2Health 状态机", 5.25, 3.22, 2.8, 0.3, 13, PURPLE, True, PP_ALIGN.CENTER)
    priorities = [
        ("1", "地面备份中心", BLUE, PALE_BLUE),
        ("2", "高空侦察节点", TEAL, PALE_TEAL),
        ("3", "资源集群代表", GREEN, PALE_GREEN),
        ("4", "拍卖/CBBA 保底", AMBER, PALE_AMBER),
    ]
    for idx, (n, text, color, fill) in enumerate(priorities):
        x = 1.0 + idx * 3.0
        add_rect(slide, x, 4.1, 2.55, 1.0, fill, color)
        add_text(slide, n, x + 0.12, 4.25, 0.4, 0.36, 20, color, True, PP_ALIGN.CENTER)
        add_text(slide, text, x + 0.58, 4.28, 1.75, 0.3, 12, INK, True, PP_ALIGN.CENTER)
    add_card(slide, "主动降级输入", "D1/D2 不确定性、D3 计划新鲜度、D5 视觉一致性、通信与二级 readiness。", 0.95, 5.72, 5.45, 0.82, AMBER, PALE_AMBER, 12.5, 10.5)
    add_card(slide, "恢复原则", "中心恢复后先双轨合并和 digest 校验，确认稳定后再恢复主控。", 6.75, 5.72, 5.55, 0.82, TEAL, PALE_TEAL, 12.5, 10.5)
    add_footer(slide, 15)

    # 16 D4 evidence
    slide = new_slide(prs, "D4 验证结果", "二级接管、完全分布式和缺 ACK 负例均已形成结构化证据", "MODULE D4")
    add_picture_contain(slide, ASSETS["d4_curve"], 0.65, 1.72, 6.4, 3.65)
    add_rect(slide, 7.35, 1.75, 5.3, 3.55, WHITE, LINE)
    add_text(slide, "联盟状态验收", 7.65, 1.98, 2.0, 0.35, 15, PURPLE, True)
    add_status_dot(slide, 7.75, 2.65, TEAL, "二级节点接管", "ACK 3/3")
    add_status_dot(slide, 7.75, 3.35, GREEN, "完全分布式 peer", "ACK 3/3")
    add_status_dot(slide, 7.75, 4.05, RED, "缺 ACK 场景", "2/3 → aborted")
    add_text(slide, "缺 ACK 时三个 T001 成员 hold_for_review，D7 许可为 0。", 7.75, 4.62, 4.45, 0.38, 11, MUTED)
    add_card(slide, "已完成", "commit / ACK / epoch / lease 正负例，secondary 与 peer 均可成为计划 owner。", 0.8, 5.62, 3.65, 0.85, GREEN, PALE_GREEN, 12.5, 10)
    add_card(slide, "仍需", "成员退出与补位、分区恢复、中心恢复 digest、误降级率的完整扰动矩阵。", 4.65, 5.62, 3.65, 0.85, AMBER, PALE_AMBER, 12.5, 10)
    add_card(slide, "P2 边界", "MIT/CA-CBBA 未接入；当前仍为本地轻量 CBBA 与原子 commit 主线。", 8.5, 5.62, 3.65, 0.85, PURPLE, PALE_PURPLE, 12.5, 10)
    add_footer(slide, 16)

    # 17 D5 geometry
    slide = new_slide(prs, "D5 终端视觉配准与身份认证", "相机看到的最近目标，不一定是中心分配的目标", "MODULE D5")
    add_module_banner(slide, "D5", "末端配准", "GlobalTrack → 图像 → 锁定", color=RED, fill=PALE_RED)
    add_node(slide, "中心 GlobalTrack", "预测位置 + 协方差\n当前 assignment", 0.85, 3.35, 2.2, 1.05, BLUE, PALE_BLUE)
    add_arrow(slide, 3.1, 3.88, 4.0, 3.88, MUTED, 1.4)
    add_node(slide, "相机几何投影", "K / R / t / distortion\n协方差传播", 4.05, 3.35, 2.3, 1.05, TEAL, PALE_TEAL)
    add_arrow(slide, 6.4, 3.88, 7.3, 3.88, MUTED, 1.4)
    add_node(slide, "局部 MOT", "YOLO / ByteTrack / BoT-SORT\nLocalVisualTrack", 7.35, 3.35, 2.35, 1.05, AMBER, PALE_AMBER)
    add_arrow(slide, 9.75, 3.88, 10.45, 3.88, MUTED, 1.4)
    add_node(slide, "保守决策", "locked / ambiguous\nhold / reacquire", 10.5, 3.35, 2.0, 1.05, RED, PALE_RED)
    add_card(slide, "匹配代价", "像素马氏距离 + 角速度一致性 + 时间戳 + 类别 + 友方冲突。", 1.15, 5.05, 3.45, 1.15, TEAL, PALE_TEAL, 13, 10.5)
    add_card(slide, "身份原则", "Remote ID / 签名 / 视觉标签只做正向确认；未知不等于敌方。", 4.95, 5.05, 3.45, 1.15, BLUE, PALE_BLUE, 13, 10.5)
    add_card(slide, "系统纪律", "D5 不自行换绑或改写 global_track_id；冲突时上报中心/二级节点。", 8.75, 5.05, 3.45, 1.15, RED, PALE_RED, 13, 10.5)
    add_footer(slide, 17)

    # 18 D5 cross-camera
    slide = new_slide(prs, "多相机视场重叠时如何关联", "无人机1看到 1/2/3，无人机2看到 2/3/4：先局部跟踪，再通过全局几何与时间一致性注册", "MODULE D5")
    add_rect(slide, 0.75, 1.85, 4.35, 3.55, WHITE, BLUE)
    add_text(slide, "相机 A / 拦截机 1", 1.0, 2.05, 2.2, 0.35, 14, BLUE, True)
    add_rect(slide, 1.0, 2.58, 3.8, 2.25, "F8FAFC", LINE, False)
    for i, (cx, cy, label) in enumerate([(1.45, 3.0, "1"), (2.55, 3.72, "2"), (3.75, 2.9, "3")]):
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.56), Inches(0.56))
        set_fill(s, BLUE if label == "1" else TEAL)
        s.line.fill.background()
        add_text(slide, label, cx, cy + 0.1, 0.56, 0.22, 12, WHITE, True, PP_ALIGN.CENTER)
    add_rect(slide, 8.2, 1.85, 4.35, 3.55, WHITE, RED)
    add_text(slide, "相机 B / 拦截机 2", 8.45, 2.05, 2.2, 0.35, 14, RED, True)
    add_rect(slide, 8.45, 2.58, 3.8, 2.25, "F8FAFC", LINE, False)
    for cx, cy, label in [(8.95, 3.72, "2"), (10.0, 2.9, "3"), (11.1, 3.65, "4")]:
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.56), Inches(0.56))
        set_fill(s, TEAL if label in {"2", "3"} else RED)
        s.line.fill.background()
        add_text(slide, label, cx, cy + 0.1, 0.56, 0.22, 12, WHITE, True, PP_ALIGN.CENTER)
    add_node(slide, "跨视角注册", "GlobalTrack 投影\n时间对齐\n协方差门限\nbbox 历史", 5.55, 2.3, 2.2, 1.75, TEAL, PALE_TEAL)
    add_arrow(slide, 5.1, 3.42, 5.5, 3.2, BLUE, 1.2)
    add_arrow(slide, 8.15, 3.42, 7.8, 3.2, RED, 1.2)
    add_badge(slide, "目标2：A↔B", 4.5, 4.48, 1.7, TEAL, PALE_TEAL)
    add_badge(slide, "目标3：A↔B", 6.45, 4.48, 1.7, TEAL, PALE_TEAL)
    add_card(slide, "完全无中心时", "各无人机交换带时间戳的局部轨迹摘要、bbox 尺度历史和候选集合；D4 原子联盟确认后才能执行。", 1.0, 5.72, 5.35, 0.82, PURPLE, PALE_PURPLE, 12.5, 10.5)
    add_card(slide, "不采用的捷径", "不使用 AirSim detection truth ID 在线匹配；不把本地 MOT ID 当 global_track_id；不因“最近”自动换目标。", 6.75, 5.72, 5.55, 0.82, RED, PALE_RED, 12.5, 10.5)
    add_footer(slide, 18)

    # 19 D5 results
    slide = new_slide(prs, "D5 验证结果与当前断点", "合同层已经形成，但视觉持续性仍是物理闭环首要瓶颈", "MODULE D5")
    add_picture_contain(slide, ASSETS["d5_timeline"], 0.65, 1.72, 7.2, 4.85)
    add_metric(slide, "8/10", "双 primary 视觉合同", 8.15, 1.88, 1.9, GREEN)
    add_metric(slide, "0", "错误 duplicate", 10.3, 1.88, 1.9, BLUE)
    add_card(slide, "P2 几何对照", "合成 calibration / solvePnP 后投影 RMSE：约 24.0 px → 1.63 px。", 8.15, 3.45, 4.05, 1.02, PURPLE, PALE_PURPLE, 13, 10.5)
    add_card(slide, "当前断点", "ComputerVision control_allowed=0；SimpleFlight 24/30 active pair 发生 terminal_detection_timeout。", 8.15, 4.7, 4.05, 1.02, RED, PALE_RED, 13, 10.5)
    add_text(slide, "下一步：真实图像下持续 detection、MOT 稳定性、相机外参漂移与 D7 gate 联合标定", 8.2, 5.98, 3.95, 0.5, 11, AMBER, True, PP_ALIGN.CENTER)
    add_footer(slide, 19)

    # 20 D7 guidance
    slide = new_slide(prs, "D7 比例导引：雷达中段到视觉末端", "保持经典 PN/PNG 核心公式，重点完善安全切换条件", "MODULE D7")
    add_module_banner(slide, "D7", "比例导引", "位置 PN → 视觉 PNG", color=GREEN, fill=PALE_GREEN)
    add_phase_band(slide, 0.85, 3.0, 2.2, "雷达中段 PN", "GlobalTrack 位置/速度", BLUE, PALE_BLUE)
    add_phase_band(slide, 3.45, 3.0, 2.2, "交接准备", "距离/检测框稳定", AMBER, PALE_AMBER)
    add_phase_band(slide, 6.05, 3.0, 2.2, "D3/D4/D5 门控", "身份/版本/联盟/机动", PURPLE, PALE_PURPLE)
    add_phase_band(slide, 8.65, 3.0, 2.2, "视觉 PNG", "LOS rate + TTC/VM", GREEN, PALE_GREEN)
    add_arrow(slide, 3.08, 3.45, 3.4, 3.45, MUTED, 1.4)
    add_arrow(slide, 5.68, 3.45, 6.0, 3.45, MUTED, 1.4)
    add_arrow(slide, 8.28, 3.45, 8.6, 3.45, MUTED, 1.4)
    add_card(slide, "切换必须同时满足", "分配一致、联盟 committed、D5 locked、相机质量、LOS 稳定、机动余量和剩余窗口。", 1.0, 4.55, 5.4, 1.15, GREEN, PALE_GREEN, 13, 10.5)
    add_card(slide, "任何一项不满足", "保持 radar_midcourse / hold / reacquire；不得绕过门控，不得激活 standby reserve。", 6.75, 4.55, 5.4, 1.15, RED, PALE_RED, 13, 10.5)
    add_badge(slide, "P2：3D PN", 2.1, 6.15, 1.55, BLUE, PALE_BLUE)
    add_badge(slide, "True PN", 3.9, 6.15, 1.35, TEAL, PALE_TEAL)
    add_badge(slide, "APN", 5.5, 6.15, 1.2, GREEN, PALE_GREEN)
    add_badge(slide, "FRPN 研究近似", 6.95, 6.15, 2.1, PURPLE, PALE_PURPLE)
    add_text(slide, "P2 仅离线质点对照，未替换默认位置 PN / 视觉 PNG", 9.25, 6.18, 3.15, 0.28, 10, MUTED, True, PP_ALIGN.CENTER)
    add_footer(slide, 20)

    # 21 Full-flow simulation
    slide = new_slide(prs, "全流程仿真已经形成可回放闭环", "质点模型用于算法闭环，AirSim 用于接口、视觉和 SimpleFlight 控制验证", "SIMULATION")
    add_picture_contain(slide, ASSETS["global_gif"], 0.65, 1.72, 6.0, 4.9)
    add_picture_contain(slide, ASSETS["air_3d"], 6.9, 1.72, 5.78, 4.9)
    add_badge(slide, "左：2D 全局流程 GIF", 1.9, 6.62, 2.7, TEAL, PALE_TEAL)
    add_badge(slide, "右：AirSim 受控 3D 轨迹", 8.4, 6.62, 2.85, BLUE, PALE_BLUE)
    add_footer(slide, 21)

    # 22 D6
    slide = new_slide(prs, "D6 系统级评估指标体系", "检测、跟踪、分配、降级、末端视觉、安全和物理结果统一进入评估", "MODULE D6")
    add_module_banner(slide, "D6", "评估指标", "日志消费与批量报告", color=INK, fill=WHITE)
    categories = [
        ("探测", "Pd / FAR / 漏检", BLUE, PALE_BLUE),
        ("跟踪", "RMSE / continuity / IDSW", TEAL, PALE_TEAL),
        ("分配", "重复 / 未分配高威胁", AMBER, PALE_AMBER),
        ("降级", "接管时间 / ACK / 共识", PURPLE, PALE_PURPLE),
        ("末端", "锁定 / 模糊 / hold", RED, PALE_RED),
        ("安全", "约束 / 人工覆盖", GREEN, PALE_GREEN),
    ]
    for idx, (name, body, color, fill) in enumerate(categories):
        x = 0.85 + (idx % 3) * 4.15
        y = 3.0 + (idx // 3) * 1.28
        add_card(slide, name, body, x, y, 3.65, 0.95, color, fill, 14, 10.5)
    add_rect(slide, 1.05, 5.78, 11.1, 0.72, WHITE, LINE)
    add_text(slide, "四层结果必须分开", 1.3, 5.98, 2.0, 0.28, 13, RED, True)
    add_badge(slide, "contract_allowed", 3.45, 5.97, 1.9, BLUE, PALE_BLUE, 9.2)
    add_badge(slide, "control_allowed", 5.6, 5.97, 1.8, TEAL, PALE_TEAL, 9.2)
    add_badge(slide, "mode_switched", 7.65, 5.97, 1.65, AMBER, PALE_AMBER, 9.2)
    add_badge(slide, "physical_intercept", 9.55, 5.97, 2.15, GREEN, PALE_GREEN, 9.2)
    add_text(slide, "D6 只消费日志，不参与控制。", 4.3, 6.7, 4.8, 0.24, 10.5, MUTED, True, PP_ALIGN.CENTER)
    add_footer(slide, 22)

    # 23 validation
    slide = new_slide(prs, "当前 AirSim 验证结论", "合同层达标，物理闭环尚未达标：两者必须分开汇报", "VALIDATION")
    add_picture_contain(slide, ASSETS["cv_accept"], 0.65, 1.72, 6.0, 3.0)
    add_picture_contain(slide, ASSETS["sf_diag"], 6.85, 1.72, 5.85, 3.0)
    add_card(slide, "ComputerVision 合同层", "双 primary 8/10；IDSW=0；错误重复锁=0；二级/peer ACK 3/3；缺 ACK 正确 aborted。", 0.75, 5.0, 3.75, 1.15, GREEN, PALE_GREEN, 13, 10.5)
    add_card(slide, "SimpleFlight 物理层", "10 seeds / 30 active pair / 0 命中；24 次末端检测超时；平均最小距离 29.849 m。", 4.8, 5.0, 3.75, 1.15, RED, PALE_RED, 13, 10.5)
    add_card(slide, "正确解释", "CV 证明合同与安全门控；15 秒、2 Hz 只定位断点，不能评价 PN/PNG 优劣。", 8.85, 5.0, 3.75, 1.15, AMBER, PALE_AMBER, 13, 10.5)
    add_text(slide, "下一轮优先：相机取景 → detection 持续性 → D5 lock → D7 gate → closing speed", 1.2, 6.5, 10.9, 0.35, 13, INK, True, PP_ALIGN.CENTER)
    add_footer(slide, 23)

    # 24 maturity
    slide = new_slide(prs, "P1 / P2 成熟度与技术边界", "已完成的不夸大，未完成的可量化", "MATURITY")
    headers = ["模块", "P1 合同", "P1 物理/长期", "P2 对照"]
    col_x = [0.65, 2.25, 5.25, 9.05]
    widths = [1.4, 2.8, 3.6, 3.65]
    for x, w, label in zip(col_x, widths, headers):
        add_rect(slide, x, 1.78, w, 0.5, INK, INK, False)
        add_text(slide, label, x + 0.05, 1.91, w - 0.1, 0.22, 11, WHITE, True, PP_ALIGN.CENTER)
    rows = [
        ("D1", "完成", "长 replay / CI 标定", "未启动第三方"),
        ("D2", "完成", "真实 crossing replay", "Stone Soup/FilterPy adapter"),
        ("D3", "完成", "动态非等量标定", "OR-Tools 仅接口"),
        ("D4", "完成", "扰动矩阵 / 重构", "外部 CBBA 未接入"),
        ("D5", "8/10", "持续检测 / control gate", "PnP 合成对照"),
        ("D6", "完成", "长期场景库 / CI", "motmetrics 两帧 smoke"),
        ("D7", "完成", "0/30 物理闭环", "3D/APN/FRPN 质点"),
    ]
    for idx, row in enumerate(rows):
        y = 2.35 + idx * 0.59
        fill = WHITE if idx % 2 == 0 else "EEF1F4"
        for x, w, text in zip(col_x, widths, row):
            add_rect(slide, x, y, w, 0.52, fill, LINE, False, 0.6)
            color = GREEN if text in {"完成", "8/10"} else INK
            add_text(slide, text, x + 0.06, y + 0.12, w - 0.12, 0.25, 10.2, color, text in {"完成", "8/10"}, PP_ALIGN.CENTER)
    add_card(slide, "总体判断", "P1 合同层可支撑继续仿真集成；P1 物理闭环是下一阶段主线；P2 只做隔离对照，不替换默认算法。", 1.05, 6.62, 11.2, 0.5, BLUE, PALE_BLUE, 11.5, 9.5)
    add_footer(slide, 24)

    # 25 roadmap
    slide = new_slide(prs, "下一阶段工作计划与资源建议", "先关闭物理闭环，再进入更真实飞控和硬件联试", "ROADMAP")
    add_phase_band(slide, 0.8, 1.85, 2.75, "阶段 1：视觉漏斗", "检测持续性 / MOT / 外参", RED, PALE_RED)
    add_phase_band(slide, 3.85, 1.85, 2.75, "阶段 2：物理闭环", "90 s / 更高频 / 多 seed", AMBER, PALE_AMBER)
    add_phase_band(slide, 6.9, 1.85, 2.75, "阶段 3：鲁棒协同", "成员重构 / 网络分区 / 恢复", TEAL, PALE_TEAL)
    add_phase_band(slide, 9.95, 1.85, 2.55, "阶段 4：工程化", "PX4/ROS2/硬件在环", BLUE, PALE_BLUE)
    add_arrow(slide, 3.58, 2.28, 3.82, 2.28, MUTED, 1.5)
    add_arrow(slide, 6.63, 2.28, 6.87, 2.28, MUTED, 1.5)
    add_arrow(slide, 9.68, 2.28, 9.92, 2.28, MUTED, 1.5)
    add_text(slide, "资源与组织建议", 0.9, 3.35, 1.8, 0.35, 16, INK, True)
    decisions = [
        ("资源聚焦", "短期停止扩展新算法名录，集中攻克 D5/D7 物理闭环。", RED, PALE_RED),
        ("验证门槛", "以长时多 seed、分层失败原因和可复现报告作为阶段验收。", GREEN, PALE_GREEN),
        ("平台路线", "AirSim 继续作为统一仿真入口；P2 依赖保持隔离，不污染主线。", BLUE, PALE_BLUE),
        ("工程准备", "同步冻结传感器、通信、算力和相机标定接口，为后续硬件在环做准备。", TEAL, PALE_TEAL),
    ]
    for idx, (title, body, color, fill) in enumerate(decisions):
        x = 0.9 + (idx % 2) * 6.0
        y = 3.9 + (idx // 2) * 1.18
        add_card(slide, title, body, x, y, 5.55, 0.95, color, fill, 13, 10.5)
    add_rect(slide, 1.05, 6.38, 11.2, 0.58, DARK, DARK)
    add_text(slide, "阶段验收目标：物理闭环完成，失败原因可分层统计，实验结果可复现", 1.3, 6.53, 10.7, 0.28, 15, WHITE, True, PP_ALIGN.CENTER)
    add_footer(slide, 25, "MSM C-UAS · 领导汇报")

    # 26 Appendix: node capability targets
    slide = new_slide(prs, "附录 A：主要节点能力指标建议", "指标来自当前 D1-D7 处理周期、AirSim 验证条件和工程样机需求分解", "APPENDIX")
    node_specs = [
        (
            "中心节点 C2",
            ["CPU：16-32 核", "内存：64-128 GB", "存储：4-8 TB NVMe", "GPU：NVIDIA L4 / RTX 6000 Ada 级", "航迹到新计划：<300 ms"],
            BLUE,
            PALE_BLUE,
        ),
        (
            "地面雷达",
            ["更新率：10 Hz 或更高", "角精度：约 0.5 deg 级", "距离误差：1-6 m 量级", "端到端延时：<100-200 ms", "必须输出时间戳、协方差和质量状态"],
            TEAL,
            PALE_TEAL,
        ),
        (
            "高空侦察二级节点",
            ["试验高度：50-200 m 分档", "续航：机动平台 30-60 min；系留备选 24-50 h", "载荷：2-5 kg", "EO/IR：4K/40MP + 热成像 + 云台/LRF", "数据链：MANET；系留备选可提供 100 Mbps"],
            PURPLE,
            PALE_PURPLE,
        ),
        (
            "拦截无人机",
            ["速度比：拦截机/目标 >=1.5-2.0", "相机：2-4 MP 全局快门、60 fps", "机载算力：Orin Nano/NX 级", "续航：10-20 min 热备/巡逻", "热备起飞响应：2-10 s"],
            GREEN,
            PALE_GREEN,
        ),
    ]
    for idx, (title, items, color, fill) in enumerate(node_specs):
        x = 0.75 + (idx % 2) * 6.2
        y = 1.72 + (idx // 2) * 2.55
        add_rect(slide, x, y, 5.85, 2.2, WHITE, color, False, 1.3)
        add_rect(slide, x, y, 5.85, 0.48, fill, color, False, 0.8)
        add_text(slide, title, x + 0.18, y + 0.12, 5.45, 0.25, 14, color, True)
        add_bullets(slide, items, x + 0.2, y + 0.62, 5.35, 1.35, 10.5, INK, 1.0)
    add_rect(slide, 0.75, 6.87, 12.0, 0.22, "F2F4F6", LINE, False)
    add_text(slide, "工程含义：D1-D4 的数学计算量较小，主要资源消耗集中在 D5 多路视频/MOT、二级视频 cue 和 D6 批量日志。", 1.05, 6.75, 11.4, 0.3, 10.5, MUTED, True, PP_ALIGN.CENTER)
    add_footer(slide, 26)

    # 27 Appendix: communication and timing budgets
    slide = new_slide(prs, "附录 B：通信带宽与响应时间预算", "默认传输结构化元数据；视频只在指定区域、指定资源之间定向传输", "APPENDIX")
    headers = ["链路/模块", "结构化数据预算", "工程建议", "时延/响应要求"]
    col_x = [0.65, 3.0, 5.6, 9.25]
    widths = [2.25, 2.5, 3.55, 3.45]
    for x, w, label in zip(col_x, widths, headers):
        add_rect(slide, x, 1.75, w, 0.52, DARK, DARK, False)
        add_text(slide, label, x + 0.06, 1.89, w - 0.12, 0.24, 10.5, WHITE, True, PP_ALIGN.CENTER)
    link_rows = [
        ("C2—雷达", "<10 Mbps", "1 GbE/专线 + 时间同步", "雷达端到端 <100-200 ms"),
        ("C2—二级节点", "1-10 Mbps metadata", "系留 100 Mbps + 无线备份", "heartbeat stale 约 2 s"),
        ("二级—拦截机", "1-10 Mbps metadata", "10-100 Mbps MANET；视频定向", "跨相机 skew <50-100 ms"),
        ("C2—拦截机", "0.5-2 Mbps/架", "控制/状态 QoS 高于视频", "计划与 bbox <100-200 ms"),
        ("拦截机—拦截机", "0.1-1 Mbps/架", "低时延 MANET；丢包可观测", "CBBA round 约 0.5 s"),
        ("D7 单 pair", "20-50 kB/s", "只传 Track/Plan/bbox，不传 PNG", "单 pair 计算 <5 ms；上限 0.35 s"),
        ("视频 cue / replay", "5-50 Mbps/路", "局部回传或离线保存", "不进入控制安全链路"),
    ]
    for idx, row in enumerate(link_rows):
        y = 2.29 + idx * 0.57
        fill = "F2F4F6" if idx % 2 == 0 else WHITE
        for x, w, text in zip(col_x, widths, row):
            add_rect(slide, x, y, w, 0.52, fill, LINE, False, 0.5)
            add_text(slide, text, x + 0.07, y + 0.13, w - 0.14, 0.25, 9.8, INK, idx == 0 and x == col_x[0], PP_ALIGN.CENTER)
    add_rect(slide, 0.8, 6.48, 11.7, 0.55, PALE_AMBER, AMBER, False)
    add_text(slide, "带宽设计原则：控制与安全状态优先，态势元数据持续传输，多路视频不得默认全网广播。", 1.05, 6.63, 11.15, 0.28, 12.5, INK, True, PP_ALIGN.CENTER)
    add_footer(slide, 27)

    # 28 Appendix: component selection and budgetary estimate
    slide = new_slide(prs, "附录 C：关键部件选型与预算级估算", "参考公开产品能力与市场资料；具体型号、价格和合规条件需在样机阶段询价确认", "APPENDIX")
    headers = ["系统层", "建议配置/参考产品", "预算级成本"]
    col_x = [0.65, 2.65, 9.75]
    widths = [1.9, 7.0, 2.9]
    for x, w, label in zip(col_x, widths, headers):
        add_rect(slide, x, 1.75, w, 0.52, DARK, DARK, False)
        add_text(slide, label, x + 0.06, 1.89, w - 0.12, 0.24, 10.5, WHITE, True, PP_ALIGN.CENTER)
    component_rows = [
        ("中心 C2", "16-32 核、64-128 GB、4-8 TB NVMe、L4/RTX 6000 Ada 级 GPU", "USD 15k-50k"),
        ("主雷达", "EchoShield/同级 4D C-UAS 雷达；10 Hz、0.5 deg、带协方差航迹", "USD 80k-250k/扇区"),
        ("二级侦察平台", "机动高空平台为当前主线；Elistair Orion 2.2 TE/同级作为长期值守备选", "USD 25k-300k/套"),
        ("二级 EO/IR", "DJI H30T 级或更高端 ISR 吊舱；高分辨率 EO、热成像、LRF", "USD 10k-100k+"),
        ("拦截无人机", "Pixhawk 6X/PX4 + Orin Nano/NX + 全局快门相机 + MANET", "USD 8k-25k+/架"),
        ("MANET 节点", "Doodle Labs / Silvus / Microhard 等；metadata 优先，视频定向", "USD 5k-20k+/节点"),
        ("时间同步/标定", "PTP/GNSS disciplined clock；相机内外参与多节点时间基准", "USD 2k-10k"),
    ]
    for idx, row in enumerate(component_rows):
        y = 2.29 + idx * 0.57
        fill = "F2F4F6" if idx % 2 == 0 else WHITE
        for x, w, text in zip(col_x, widths, row):
            add_rect(slide, x, y, w, 0.52, fill, LINE, False, 0.5)
            add_text(slide, text, x + 0.07, y + 0.12, w - 0.14, 0.27, 9.5, INK, x == col_x[0], PP_ALIGN.CENTER)
    add_rect(slide, 0.8, 6.45, 7.6, 0.58, PALE_BLUE, BLUE, False)
    add_text(slide, "5v5 工程样机粗预算：USD 300k-900k+", 1.05, 6.6, 7.1, 0.3, 14, BLUE, True, PP_ALIGN.CENTER)
    add_rect(slide, 8.65, 6.45, 3.85, 0.58, PALE_RED, RED, False)
    add_text(slide, "不含场地、合规、保险、维护和人员", 8.85, 6.6, 3.45, 0.3, 10.5, RED, True, PP_ALIGN.CENTER)
    add_footer(slide, 28)

    assert len(prs.slides) == 28, f"expected 28 slides, got {len(prs.slides)}"
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = create_deck()
    print(path)
